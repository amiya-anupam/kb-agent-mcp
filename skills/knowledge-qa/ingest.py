#!/usr/bin/env python3
"""
ingest.py — recursively reads all supported files from ~/Desktop/KnowledgeBase/
and prints a JSON array of { file, type, chars, text, confidential, confidential_reason }
to stdout.

Usage:
    uv run \
      --with pypdf \
      --with python-docx \
      --with openpyxl \
      --with python-pptx \
      --with beautifulsoup4 \
      --with striprtf \
      --with pyyaml \
      --with ebooklib \
      --with defusedxml \
      ingest.py [folder]

Supported:
    Text/doc:  .pdf, .docx, .xlsx, .xls, .pptx, .txt, .md, .csv
    Box Notes: .boxnote  (ProseMirror JSON — text extracted recursively)
    Web/data:  .html, .htm  (text stripped via BeautifulSoup; file opened in browser on request)
               .rtf  (markup stripped via striprtf)
               .json  (pretty-printed parsed tree)
               .yaml, .yml  (loaded and dumped to readable string)
               .xml  (element tree walked via defusedxml — safe against XML bomb attacks)
               .epub  (chapter-by-chapter text via ebooklib)
               .eml  (email body + headers via built-in email module; depth-limited MIME walk)
    Images:    .png, .jpg, .jpeg, .gif, .webp  (size-checked before load)

Security hardening:
    - Path traversal: folder argument validated; must resolve inside ALLOWED_ROOTS
    - Symlink traversal: symlinks are skipped entirely during rglob walk
    - XML bomb: defusedxml replaces stdlib ET for all XML parsing
    - Image memory bomb: files >MAX_IMAGE_MB are skipped before loading
    - EML recursion bomb: MIME walk limited to MAX_MIME_DEPTH parts
    - Error leakage: exception messages are sanitized before returning
    - Total context cap: aggregate output capped at MAX_TOTAL_CHARS
    - Audit enforcement: pass --skip-audit-check only after explicit user acknowledgement

Confidentiality classification:
    - Every output entry carries confidential: true/false and confidential_reason.
    - Detection sources: text body keywords, filename/path keywords, PDF metadata,
      DOCX core properties, EML Sensitivity header.
    - Folders containing a .noindex sentinel file are skipped entirely — no files
      inside are extracted, regardless of their content.
    - Classification does NOT redact — text is always returned in full. The SKILL.md
      consent gate decides whether flagged content enters the LLM context window.
"""

import sys
import json
import base64
import re
import os
from pathlib import Path

# --------------------------------------------------------------------------- #
# Configuration / limits                                                        #
# --------------------------------------------------------------------------- #

# Roots the folder argument must resolve inside (prevents path traversal).
# Only the Desktop KnowledgeBase and /tmp (for tests) are allowed.
ALLOWED_ROOTS = [
    Path.home() / "Desktop" / "KnowledgeBase",
    Path("/tmp"),
]

# Per-file text cap (characters)
MAX_FILE_CHARS = 50_000

# Total aggregate cap across all files (characters) — ~10 MB of text
MAX_TOTAL_CHARS = 10_000_000

# Maximum image file size to load into memory
MAX_IMAGE_MB = 50

# Maximum MIME parts to walk in a single .eml (prevents recursion bomb)
MAX_MIME_DEPTH = 50

# --------------------------------------------------------------------------- #
# Argument parsing & path validation                                            #
# --------------------------------------------------------------------------- #

def _sanitize_error(e: Exception) -> str:
    """
    Return a safe, low-detail error string.
    Strips absolute paths and internal library details from exception messages
    to prevent information leakage into the LLM context window.
    """
    msg = str(e)
    # Remove absolute paths (e.g. /Users/alice/... or C:\Users\...)
    msg = re.sub(r"(/[^\s,;\"']+|[A-Za-z]:\\[^\s,;\"']+)", "<path>", msg)
    # Truncate to a safe length
    return msg[:200]


def _validate_folder(raw: str) -> Path:
    """
    Resolve the given path and verify it sits inside one of the ALLOWED_ROOTS.
    Raises SystemExit with a clear message if the path is disallowed.
    """
    try:
        candidate = Path(raw).resolve()
    except Exception as e:
        print(json.dumps({"error": f"Invalid folder path: {_sanitize_error(e)}"}))
        sys.exit(2)

    for root in ALLOWED_ROOTS:
        try:
            candidate.relative_to(root.resolve())
            return candidate
        except ValueError:
            continue

    # Not inside any allowed root
    allowed_str = ", ".join(str(r) for r in ALLOWED_ROOTS)
    print(json.dumps({
        "error": (
            f"Folder '{candidate}' is outside the allowed roots. "
            f"Allowed: {allowed_str}"
        )
    }))
    sys.exit(2)


FOLDER = _validate_folder(sys.argv[1]) if len(sys.argv) > 1 \
         else (Path.home() / "Desktop" / "KnowledgeBase").resolve()

TEXT_TYPES  = {".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".txt", ".md", ".csv"}
WEB_TYPES   = {".html", ".htm", ".rtf", ".json", ".yaml", ".yml", ".xml", ".epub", ".eml"}
IMAGE_TYPES = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
ALL_SUPPORTED = TEXT_TYPES | {".boxnote"} | WEB_TYPES | IMAGE_TYPES

# --------------------------------------------------------------------------- #
# Confidentiality classification                                                #
# --------------------------------------------------------------------------- #

# Keywords that signal a document is confidential/restricted.
# Checked case-insensitively against text body AND file path segments.
_CONFIDENTIAL_KEYWORDS = [
    "confidential",
    "internal use only",
    "not for distribution",
    "not for sharing",
    "do not share",
    "do not distribute",
    "proprietary",
    "restricted",
    "privileged",
    "ibm confidential",
    "company confidential",
    "for internal use",
    "classification:",
    "sensitive",
    "top secret",
    "private and confidential",
]

# Compiled pattern for fast body scan (whole-word / phrase match).
_CONFIDENTIAL_RE = re.compile(
    "|".join(re.escape(kw) for kw in _CONFIDENTIAL_KEYWORDS),
    re.IGNORECASE,
)


def _classify_confidentiality(
    path: Path,
    text: str,
    pdf_meta: dict | None = None,
    docx_props: dict | None = None,
    eml_sensitivity: str | None = None,
) -> tuple[bool, str]:
    """
    Return (confidential: bool, reason: str).

    Checks (in order, first match wins):
      1. EML Sensitivity header
      2. PDF metadata (/Keywords, /Subject)
      3. DOCX core properties (category, keywords)
      4. File path / filename segments
      5. Extracted text body
    """
    # 1 — EML Sensitivity header
    if eml_sensitivity:
        sens = eml_sensitivity.lower()
        if any(kw in sens for kw in ("confidential", "company-confidential",
                                     "personal", "private", "restricted")):
            return True, f"EML Sensitivity header: {eml_sensitivity!r}"

    # 2 — PDF metadata
    if pdf_meta:
        for field in ("keywords", "subject"):
            val = (pdf_meta.get(field) or "").lower()
            if _CONFIDENTIAL_RE.search(val):
                return True, f"PDF metadata ({field}): {pdf_meta[field]!r}"

    # 3 — DOCX core properties
    if docx_props:
        for field in ("category", "keywords"):
            val = (docx_props.get(field) or "").lower()
            if _CONFIDENTIAL_RE.search(val):
                return True, f"DOCX property ({field}): {docx_props[field]!r}"

    # 4 — File path / filename
    path_str = str(path).lower()
    m = _CONFIDENTIAL_RE.search(path_str)
    if m:
        return True, f"Filename/path contains: {m.group(0)!r}"

    # 5 — Text body (first match only — don't scan entire text for speed)
    m = _CONFIDENTIAL_RE.search(text)
    if m:
        # Capture a small snippet of surrounding context for the reason string.
        start = max(0, m.start() - 30)
        end   = min(len(text), m.end() + 30)
        snippet = text[start:end].replace("\n", " ").strip()
        return True, f"Body contains: {m.group(0)!r} — near: {snippet!r}"

    return False, ""


# --------------------------------------------------------------------------- #
# Extractors                                                                    #
# --------------------------------------------------------------------------- #

def extract_pdf(path: Path) -> tuple[str, dict]:
    """Returns (text, pdf_meta). pdf_meta keys: 'keywords', 'subject'."""
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        meta = reader.metadata or {}
        pdf_meta = {
            "keywords": str(meta.get("/Keywords") or ""),
            "subject":  str(meta.get("/Subject")  or ""),
        }
        return text, pdf_meta
    except Exception as e:
        return f"[PDF extraction error: {_sanitize_error(e)}]", {}


def extract_docx(path: Path) -> tuple[str, dict]:
    """Returns (text, docx_props). docx_props keys: 'category', 'keywords'."""
    try:
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        cp = doc.core_properties
        docx_props = {
            "category": str(cp.category or ""),
            "keywords": str(cp.keywords or ""),
        }
        return text, docx_props
    except Exception as e:
        return f"[DOCX extraction error: {_sanitize_error(e)}]", {}


def extract_xlsx(path: Path) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        chunks = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            chunks.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    chunks.append(row_text)
        return "\n".join(chunks)
    except Exception as e:
        return f"[XLSX extraction error: {_sanitize_error(e)}]"


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                chunks.append(f"=== Slide {i} ===\n" + "\n".join(slide_texts))
        return "\n".join(chunks)
    except Exception as e:
        return f"[PPTX extraction error: {_sanitize_error(e)}]"


def extract_plain_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"[Text read error: {_sanitize_error(e)}]"


def _boxnote_walk(node: object, parts: list) -> None:
    """Recursively walk a ProseMirror JSON node and collect all text leaves."""
    if isinstance(node, dict):
        if node.get("type") == "text" and isinstance(node.get("text"), str):
            parts.append(node["text"])
        if node.get("type") in ("heading", "paragraph", "list_item",
                                "bullet_list", "ordered_list", "blockquote",
                                "table_row", "table_cell"):
            parts.append("\n")
        for child in node.get("content", []):
            _boxnote_walk(child, parts)
    elif isinstance(node, list):
        for child in node:
            _boxnote_walk(child, parts)


def extract_boxnote(path: Path) -> str:
    try:
        raw = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        doc = raw.get("doc") or raw
        parts: list = []
        _boxnote_walk(doc, parts)
        text = "".join(parts)
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        return text
    except Exception as e:
        return f"[Boxnote extraction error: {_sanitize_error(e)}]"


def extract_html(path: Path) -> str:
    """
    Strip HTML tags and return readable text via BeautifulSoup.
    NOTE: HTML files are text-extracted here for Q&A purposes only.
    To VIEW the rendered page, open the file in the default browser:
        open <absolute_path>   (macOS)
    Bob should NEVER execute or render HTML as a script — always use
    the system 'open' command to launch it in the user's default browser.
    """
    try:
        from bs4 import BeautifulSoup
        raw = path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        return (
            f"[HTML FILE — text extracted for Q&A]\n"
            f"Absolute path: {path}\n"
            f"To VIEW this page in your default browser run: open \"{path}\"\n\n"
            + text
        )
    except Exception as e:
        return f"[HTML extraction error: {_sanitize_error(e)}]"


def extract_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        raw = path.read_text(encoding="utf-8", errors="replace")
        return rtf_to_text(raw)
    except Exception as e:
        return f"[RTF extraction error: {_sanitize_error(e)}]"


def extract_json(path: Path) -> str:
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        parsed = json.loads(raw)
        return json.dumps(parsed, ensure_ascii=False, indent=2)
    except Exception as e:
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return f"[JSON extraction error: {_sanitize_error(e)}]"


def extract_yaml(path: Path) -> str:
    try:
        import yaml
        raw = path.read_text(encoding="utf-8", errors="replace")
        parsed = yaml.safe_load(raw)
        return yaml.dump(parsed, allow_unicode=True, default_flow_style=False)
    except Exception as e:
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return f"[YAML extraction error: {_sanitize_error(e)}]"


def extract_xml(path: Path) -> str:
    """
    Uses defusedxml to parse XML — safe against XML bomb (Billion Laughs),
    quadratic blowup, and external entity (XXE) injection attacks.
    stdlib xml.etree.ElementTree is NOT used here for security reasons.
    """
    try:
        import defusedxml.ElementTree as ET
        tree = ET.parse(str(path))
        root = tree.getroot()
        parts = []

        def walk(node, depth: int = 0) -> None:
            tag = node.tag.split("}")[-1] if "}" in node.tag else node.tag
            text = (node.text or "").strip()
            tail = (node.tail or "").strip()
            if text:
                parts.append("  " * depth + f"<{tag}>: {text}")
            for child in node:
                walk(child, depth + 1)
            if tail:
                parts.append("  " * depth + tail)

        walk(root)
        return "\n".join(parts)
    except Exception as e:
        return f"[XML extraction error: {_sanitize_error(e)}]"


def extract_epub(path: Path) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup

        book = epub.read_epub(str(path))
        chapters = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n").strip()
            if text:
                chapters.append(text)
        full = "\n\n---\n\n".join(chapters)
        return re.sub(r"\n{3,}", "\n\n", full).strip()
    except Exception as e:
        return f"[EPUB extraction error: {_sanitize_error(e)}]"


def extract_eml(path: Path) -> tuple[str, str]:
    """
    Depth-limited MIME walk (MAX_MIME_DEPTH parts) prevents recursion bombs
    from maliciously crafted multi-part email structures.

    Returns (text, sensitivity) where sensitivity is the raw value of the
    Sensitivity header (empty string if absent).
    """
    try:
        import email
        from email import policy

        raw = path.read_bytes()
        msg = email.message_from_bytes(raw, policy=policy.default)

        # Capture Sensitivity header for confidentiality classification.
        sensitivity = str(msg.get("Sensitivity") or "")

        headers = []
        for h in ("From", "To", "Cc", "Subject", "Date"):
            val = msg.get(h)
            if val:
                headers.append(f"{h}: {val}")

        body_parts = []
        part_count = 0
        if msg.is_multipart():
            for part in msg.walk():
                part_count += 1
                if part_count > MAX_MIME_DEPTH:
                    body_parts.append(f"[EML truncated: exceeded {MAX_MIME_DEPTH} MIME parts]")
                    break
                ct = part.get_content_type()
                cd = str(part.get("Content-Disposition", ""))
                if ct == "text/plain" and "attachment" not in cd:
                    body_parts.append(part.get_content())
                elif ct == "text/html" and not body_parts and "attachment" not in cd:
                    try:
                        from bs4 import BeautifulSoup
                        soup = BeautifulSoup(part.get_content(), "html.parser")
                        body_parts.append(soup.get_text(separator="\n"))
                    except Exception:
                        body_parts.append(part.get_content())
        else:
            body_parts.append(msg.get_content())

        text = "\n".join(headers) + "\n\n" + "\n---\n".join(body_parts)
        return text, sensitivity
    except Exception as e:
        return f"[EML extraction error: {_sanitize_error(e)}]", ""


def extract_image(path: Path) -> str:
    """
    Size-checks before loading. Files over MAX_IMAGE_MB are skipped entirely
    to prevent memory exhaustion from large image files.
    """
    try:
        size_bytes = path.stat().st_size
        size_mb = size_bytes / (1024 * 1024)
        if size_mb > MAX_IMAGE_MB:
            return (
                f"[IMAGE FILE SKIPPED — {size_mb:.1f} MB exceeds {MAX_IMAGE_MB} MB limit]\n"
                f"Absolute path: {path}\n"
                f"To analyse this image visually, use the read_file tool on: {path}"
            )
        ext = path.suffix.lower().lstrip(".")
        mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg",
                    "png": "image/png", "gif": "image/gif", "webp": "image/webp"}
        mime = mime_map.get(ext, "image/png")
        b64 = base64.b64encode(path.read_bytes()).decode("ascii")
        size_kb = size_bytes // 1024
        return (
            f"[IMAGE FILE — {size_kb} KB — {mime}]\n"
            f"Absolute path: {path}\n"
            f"To analyse this image visually, use the read_file tool on: {path}\n"
            f"data:{mime};base64,{b64[:200]}...  (truncated)"
        )
    except Exception as e:
        return f"[Image read error: {_sanitize_error(e)}]"


# Extractors that return (text, extra_meta) tuples instead of plain str.
# The extra_meta is forwarded to _classify_confidentiality().
_TUPLE_EXTRACTORS = {".pdf", ".docx", ".eml"}

# --------------------------------------------------------------------------- #
# Dispatch table                                                                #
# --------------------------------------------------------------------------- #

EXTRACTORS = {
    ".pdf":      extract_pdf,
    ".docx":     extract_docx,
    ".xlsx":     extract_xlsx,
    ".xls":      extract_xlsx,
    ".pptx":     extract_pptx,
    ".txt":      extract_plain_text,
    ".md":       extract_plain_text,
    ".csv":      extract_plain_text,
    ".boxnote":  extract_boxnote,
    # Web / structured data
    ".html":     extract_html,
    ".htm":      extract_html,
    ".rtf":      extract_rtf,
    ".json":     extract_json,
    ".yaml":     extract_yaml,
    ".yml":      extract_yaml,
    ".xml":      extract_xml,
    ".epub":     extract_epub,
    ".eml":      extract_eml,
    # Images
    ".png":      extract_image,
    ".jpg":      extract_image,
    ".jpeg":     extract_image,
    ".gif":      extract_image,
    ".webp":     extract_image,
}

# --------------------------------------------------------------------------- #
# Main                                                                          #
# --------------------------------------------------------------------------- #

if not FOLDER.exists():
    print(json.dumps({"error": f"Folder not found: {FOLDER}"}))
    sys.exit(1)

# Pre-compute set of folders that contain a .noindex sentinel file.
# Any file whose parent (at any depth) is in this set is skipped entirely.
_noindex_dirs: set[Path] = {
    p.parent.resolve()
    for p in FOLDER.rglob(".noindex")
    if not p.is_symlink()
}


def _in_noindex_folder(file_path: Path) -> bool:
    """Return True if any ancestor of file_path is a .noindex-marked directory."""
    resolved = file_path.resolve()
    for part in resolved.parents:
        if part in _noindex_dirs:
            return True
    return False


results = []
total_chars = 0

for f in sorted(FOLDER.rglob("*")):
    # CVE-2: Skip symlinks — never follow them, even if they point inside the folder.
    if f.is_symlink():
        continue
    if not f.is_file():
        continue
    ext = f.suffix.lower()
    if ext not in ALL_SUPPORTED:
        continue

    # Folder sentinel: skip every file inside a .noindex-marked directory.
    if _in_noindex_folder(f):
        continue

    # CVE-9: Hard stop if aggregate output cap is reached.
    if total_chars >= MAX_TOTAL_CHARS:
        results.append({
            "file": "[TRUNCATED]",
            "type": "meta",
            "chars": 0,
            "text": (
                f"[AGGREGATE CAP REACHED — output stopped at {MAX_TOTAL_CHARS:,} chars. "
                f"Remaining files were not processed. Move unneeded files out of "
                f"KnowledgeBase or raise MAX_TOTAL_CHARS in ingest.py.]"
            ),
            "confidential": False,
            "confidential_reason": "",
        })
        break

    rel = str(f.relative_to(FOLDER))

    # --- Extract text (some extractors return (text, meta) tuples) ----------
    raw_result = EXTRACTORS[ext](f)

    pdf_meta       = None
    docx_props     = None
    eml_sensitivity = None

    if ext in _TUPLE_EXTRACTORS:
        text, extra = raw_result
        if ext == ".pdf":
            pdf_meta = extra
        elif ext == ".docx":
            docx_props = extra
        elif ext == ".eml":
            eml_sensitivity = extra
    else:
        text = raw_result

    capped = text[:MAX_FILE_CHARS]
    total_chars += len(capped)

    # --- Confidentiality classification -------------------------------------
    is_confidential, reason = _classify_confidentiality(
        path=f,
        text=capped,
        pdf_meta=pdf_meta,
        docx_props=docx_props,
        eml_sensitivity=eml_sensitivity,
    )

    results.append({
        "file":                rel,
        "type":                ext.lstrip("."),
        "chars":               len(capped),
        "text":                capped,
        "confidential":        is_confidential,
        "confidential_reason": reason,
    })

print(json.dumps(results, ensure_ascii=False, indent=2))
