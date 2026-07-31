"""
kb_agent_mcp/file_parser.py
───────────────────────────
Async multi-format text extractor.

All I/O is wrapped in asyncio.to_thread() so it doesn't block the event loop.

Supported formats
-----------------
.txt / .md / .csv       plain read
.docx                   XML extraction (zipfile)
.pdf                    pypdf page iteration
.pptx / .ppt            python-pptx shape text + tables + speaker notes
.xlsx / .xls            smart streaming aggregation (large files) or openpyxl
.boxnote                JSON tree walk
.png / .jpg / .jpeg     OCR via pytesseract (optional) or PIL fallback
.gif / .webp            OCR via pytesseract (optional)
others                  filename as fallback text

Public API
----------
await extract(file_path, max_chars=None)  → str
      snippet(file_path, max_chars=2000)  → str   (sync, for indexing)
INCLUDE_EXTS                              — set of supported extensions
should_skip(path)                         — True if this file should be ignored
"""

from __future__ import annotations

import asyncio
import json
import logging
import pathlib
import re
import zipfile
from collections import defaultdict
from typing import Any

from kb_agent_mcp.config import cfg
from kb_agent_mcp.context_budget import get as _budget

logger = logging.getLogger(__name__)

# ── Constants ─────────────────────────────────────────────────────────────────

INCLUDE_EXTS: frozenset[str] = frozenset({
    ".pdf", ".docx", ".pptx", ".ppt", ".xlsx", ".xls",
    ".md", ".txt", ".csv", ".boxnote", ".doc",
    ".png", ".jpg", ".jpeg", ".gif", ".webp",
})

_SKIP_PATTERNS: frozenset[str] = frozenset({
    "readme", ".ds_store", "__pycache__",
})

_LARGE_XLSX_BYTES = 50 * 1024 * 1024  # 50 MB threshold for streaming aggregation

# ── XLSX aggregation config ───────────────────────────────────────────────────

_AGG_KEYWORDS: dict[str, str] = {
    "ut lvl 30 name dynamic": "Product (UT L30)",
    "ut l30 name": "Product (UT L30)",
    "ut l30": "Product (UT L30)",
    "product family name": "Product Family",
    "reporting product family": "Reporting Product Family",
    "product": "Product",
    "year": "Year",
    "quarter": "Quarter",
    "quarter in year": "Quarter In Year",
    "geography": "Geography",
    "geography name": "Geography",
    "market": "Market",
    "market name": "Market",
    "country": "Country",
    "finance family": "Finance Family",
    "revenue type": "Revenue Type",
    "reporting revenue type name": "Revenue Type",
    "on-prem or saas": "On-prem/SaaS",
    "division": "Division",
    "classification name": "Classification",
    "frozen client lifecycle name": "Client Lifecycle",
    "status": "Status",
}

_PREFERRED_NUM_COLS: list[str] = [
    "won", "total(cy cw won @ pc)", "rev act @ pc",
    "amount", "oppty value", "total",
]


# ── Skip helper ───────────────────────────────────────────────────────────────

def _has_noindex_ancestor(path: pathlib.Path) -> bool:
    """
    Return True if any ancestor directory of *path* (up to the KB root)
    contains a `.noindex` sentinel file.

    Imported here to avoid a circular import with security_gate — both modules
    need this logic, so file_parser owns the canonical implementation and
    security_gate delegates to it.
    """
    from kb_agent_mcp.config import cfg as _cfg
    kb_root = _cfg.kb_root_path
    for parent in path.parents:
        if (parent / ".noindex").exists():
            return True
        if parent == kb_root:
            break
    return False


def should_skip(path: pathlib.Path) -> bool:
    """Return True if this file should be excluded from indexing.

    Skips:
    • files matching _SKIP_PATTERNS (readme, .ds_store, __pycache__)
    • any file whose ancestor folder contains a `.noindex` sentinel file
    """
    name = path.name.lower()
    if any(pat in name for pat in _SKIP_PATTERNS):
        return True
    return _has_noindex_ancestor(path)


# ── Sync extractors (run in thread pool) ──────────────────────────────────────

def _extract_txt(path: pathlib.Path, max_chars: int) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")[:max_chars]


def _extract_docx(path: pathlib.Path, max_chars: int) -> str:
    with zipfile.ZipFile(path) as z:
        with z.open("word/document.xml") as f:
            xml = f.read().decode("utf-8", errors="ignore")
    text = re.sub(r"<[^>]+>", " ", xml)
    return " ".join(text.split())[:max_chars]


def _extract_pdf(path: pathlib.Path, max_chars: int) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    text = ""
    for page in reader.pages:
        text += (page.extract_text() or "") + "\n"
        if len(text) >= max_chars:
            break
    return text[:max_chars]


def _extract_pptx(path: pathlib.Path, max_chars: int) -> str:
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401 — imported for type hints in callers
    prs = Presentation(str(path))
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            # ── Text frames ────────────────────────────────────────────────────
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_parts.append(t)

            # ── Tables ─────────────────────────────────────────────────────────
            if shape.has_table:
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    non_empty = [c for c in cells if c]
                    if non_empty:
                        rows.append(" | ".join(non_empty))
                if rows:
                    slide_parts.append("\n".join(rows))

        # ── Speaker notes ──────────────────────────────────────────────────────
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                notes = notes_tf.text.strip()
                if notes:
                    slide_parts.append(f"[Notes: {notes}]")

        if slide_parts:
            parts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_parts))

        if sum(len(p) for p in parts) >= max_chars:
            break

    return "\n\n".join(parts)[:max_chars]


def _extract_image(path: pathlib.Path, max_chars: int) -> str:
    """
    Extract text from an image file.

    Strategy (in priority order):
    1. pytesseract OCR — if installed and KB_OCR_ENABLED=true (default).
    2. PIL/Pillow basic metadata fallback — image dimensions + filename.
    3. Filename-only last resort.
    """
    if not cfg.KB_OCR_ENABLED:
        return f"[Image: {path.name}]"

    # ── Try pytesseract ────────────────────────────────────────────────────────
    if cfg.KB_OCR_ENGINE in ("tesseract", "auto"):
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(str(path))
            text = pytesseract.image_to_string(img).strip()
            if text:
                return text[:max_chars]
        except ImportError:
            pass  # fall through to PIL-only fallback
        except Exception as exc:
            logger.debug("pytesseract failed for %s: %s", path.name, exc)

    # ── PIL metadata fallback ──────────────────────────────────────────────────
    try:
        from PIL import Image
        img = Image.open(str(path))
        w, h = img.size
        mode = img.mode
        return f"[Image: {path.name} | {w}×{h}px | mode={mode}]"[:max_chars]
    except ImportError:
        pass
    except Exception as exc:
        logger.debug("PIL failed for %s: %s", path.name, exc)

    return f"[Image: {path.name}]"


def _extract_boxnote(path: pathlib.Path, max_chars: int) -> str:
    """Walk a BoxNote JSON tree and collect all text nodes."""
    def _walk(node: Any, parts: list[str]) -> None:
        if isinstance(node, str):
            parts.append(node)
        elif isinstance(node, dict):
            if "text" in node and isinstance(node["text"], str):
                parts.append(node["text"])
            for v in node.values():
                _walk(v, parts)
        elif isinstance(node, list):
            for item in node:
                _walk(item, parts)

    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        parts: list[str] = []
        _walk(data, parts)
        return " ".join(parts)[:max_chars]
    except Exception:
        return f"BoxNote: {path.name}"


def _stream_xlsx_aggregate(path: pathlib.Path, max_chars: int) -> str:
    """
    Stream-aggregate a large XLSX file using raw XML parsing (iterparse).
    Never loads the full workbook into memory.
    Returns a markdown summary of totals by detected group-by dimensions.
    """
    import xml.etree.ElementTree as ET

    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            # ── Resolve sheet names ───────────────────────────────────────────
            sheet_names: list[tuple[str, str]] = []
            try:
                wb_xml = ET.fromstring(zf.read("xl/workbook.xml"))
                ns = {"w": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
                for s in wb_xml.findall(".//w:sheet", ns):
                    r_id = s.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", ""
                    )
                    sheet_names.append((s.get("name", r_id), r_id))
            except Exception as exc:
                logger.debug("Failed to parse xl/workbook.xml (%s); using default sheet name", exc)
                sheet_names = [("Sheet1", "rId1")]

            rid_to_path: dict[str, str] = {}
            try:
                rels = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
                for rel in rels:
                    rid_to_path[rel.get("Id", "")] = "xl/" + rel.get("Target", "").lstrip("/")
            except Exception as exc:
                logger.debug("Failed to parse xl/_rels/workbook.xml.rels (%s); sheet paths unavailable", exc)

            # ── Load shared strings ───────────────────────────────────────────
            shared: list[str] = []
            try:
                ss_xml = ET.fromstring(zf.read("xl/sharedStrings.xml"))
                ns_ss = {"w": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
                for si in ss_xml.findall("w:si", ns_ss):
                    parts = [t.text or "" for t in si.findall(".//w:t", ns_ss)]
                    shared.append("".join(parts))
            except Exception as exc:
                logger.debug("Failed to parse xl/sharedStrings.xml (%s); string values may be missing", exc)

            def cell_value(cell_el: Any) -> Any:
                t = cell_el.get("t", "")
                ns_v = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                v_el = cell_el.find(f"{{{ns_v}}}v")
                raw = v_el.text if v_el is not None else None
                if raw is None:
                    return None
                if t == "s":
                    try:
                        return shared[int(raw)]
                    except (IndexError, ValueError):
                        return raw
                try:
                    f = float(raw)
                    return int(f) if f == int(f) else f
                except ValueError:
                    return raw

            text_parts: list[str] = []

            for sheet_name, r_id in sheet_names:
                zip_path = rid_to_path.get(r_id, "")
                if not zip_path or zip_path not in zf.namelist():
                    for cand in [
                        "xl/worksheets/sheet1.xml",
                        f"xl/worksheets/{sheet_name}.xml",
                    ]:
                        if cand in zf.namelist():
                            zip_path = cand
                            break
                if not zip_path or zip_path not in zf.namelist():
                    continue

                text_parts.append(f"[Sheet: {sheet_name}]")

                headers: list[str] = []
                agg: dict[int, dict[str, float]] = {}
                num_ci: int | None = None
                group_cols: list[tuple[int, str]] = []
                row_count = 0

                ns_ws = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                with zf.open(zip_path) as ws_f:
                    for event, elem in ET.iterparse(ws_f, events=("end",)):
                        if elem.tag != f"{{{ns_ws}}}row":
                            continue
                        vals = [cell_value(c) for c in list(elem)]

                        if row_count == 0:
                            headers = [str(v) if v is not None else "" for v in vals]
                            h_lower = [h.lower().strip() for h in headers]

                            for pref in _PREFERRED_NUM_COLS:
                                if pref in h_lower:
                                    ci = h_lower.index(pref)
                                    num_ci = ci
                                    break

                            _agg_order = list(_AGG_KEYWORDS.keys())
                            group_cols = sorted(
                                [
                                    (ci, _AGG_KEYWORDS[h])
                                    for ci, h in enumerate(h_lower)
                                    if h in _AGG_KEYWORDS
                                ],
                                key=lambda x: _agg_order.index(
                                    next(k for k in _agg_order if _AGG_KEYWORDS[k] == x[1])
                                ),
                            )
                            _seen: set[str] = set()
                            group_cols = [
                                (ci, lbl)
                                for ci, lbl in group_cols
                                if lbl not in _seen and not _seen.add(lbl)  # type: ignore[func-returns-value]
                            ]
                            agg = {g_idx: defaultdict(float) for g_idx, _ in group_cols}
                            row_count += 1
                            elem.clear()
                            continue

                        row_count += 1
                        num_val: float | None = None
                        if num_ci is not None and num_ci < len(vals):
                            v = vals[num_ci]
                            if isinstance(v, (int, float)):
                                num_val = float(v)
                        if num_val is None:
                            for v in reversed(vals):
                                if isinstance(v, (int, float)):
                                    num_val = float(v)
                                    break

                        if num_val is not None:
                            for g_idx, _ in group_cols:
                                gval = (
                                    str(vals[g_idx])
                                    if g_idx < len(vals) and vals[g_idx] is not None
                                    else "(blank)"
                                )
                                agg[g_idx][gval] += num_val
                        elem.clear()

                rev_col = headers[num_ci] if num_ci is not None and num_ci < len(headers) else "numeric"
                text_parts.append(f"Rows: {row_count - 1}  |  Revenue column: '{rev_col}'")

                if group_cols and agg:
                    for g_idx, g_label in group_cols:
                        totals = dict(agg[g_idx])
                        if not totals:
                            continue
                        grand = sum(totals.values())
                        text_parts.append(f"--- By {g_label} ---")
                        for k, v in sorted(totals.items(), key=lambda x: -x[1])[:50]:
                            text_parts.append(f"  {k}: {v:,.2f}")
                        text_parts.append(f"  TOTAL: {grand:,.2f}\n")
                else:
                    text_parts.append(f"Columns: {', '.join(h for h in headers[:20] if h)}")
                text_parts.append("")

        return "\n".join(text_parts)[:max_chars]

    except Exception as exc:
        return f"[Large XLSX stream error: {exc}] File: {path.name}"


def _extract_xlsx(path: pathlib.Path, max_chars: int) -> str:
    """Extract text from an XLSX/XLS file. Uses streaming for large files."""
    import openpyxl

    if path.stat().st_size > _LARGE_XLSX_BYTES:
        return _stream_xlsx_aggregate(path, max_chars)

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    text_parts: list[str] = []

    sorted_sheets = sorted(
        wb.worksheets,
        key=lambda s: s.max_row or 0,
        reverse=True,
    )

    for sheet in sorted_sheets:
        row_iter = sheet.iter_rows(values_only=True)
        try:
            hdr = next(row_iter)
        except StopIteration:
            continue
        headers = [str(c) if c is not None else "" for c in hdr]
        h_lower = [h.lower().strip() for h in headers]

        probe: list = []
        for row in row_iter:
            probe.append(row)
            if len(probe) >= 200:
                break
        is_large = len(probe) >= 200

        text_parts.append(f"[Sheet: {sheet.title}]")

        if is_large:
            num_ci: int | None = None
            for pref in _PREFERRED_NUM_COLS:
                if pref in h_lower:
                    ci = h_lower.index(pref)
                    sample = [r[ci] for r in probe if ci < len(r) and r[ci] is not None]
                    if sample and all(isinstance(v, (int, float)) for v in sample):
                        num_ci = ci
                        break
            if num_ci is None:
                for ci in range(len(headers) - 1, -1, -1):
                    sample = [r[ci] for r in probe if ci < len(r) and r[ci] is not None]
                    if sample and all(isinstance(v, (int, float)) for v in sample):
                        num_ci = ci
                        break

            _agg_order = list(_AGG_KEYWORDS.keys())
            group_cols: list[tuple[int, str]] = sorted(
                [(ci, _AGG_KEYWORDS[h]) for ci, h in enumerate(h_lower) if h in _AGG_KEYWORDS],
                key=lambda x: _agg_order.index(
                    next(k for k in _agg_order if _AGG_KEYWORDS[k] == x[1])
                ),
            )
            _seen_labels: set[str] = set()
            group_cols = [
                (ci, label)
                for ci, label in group_cols
                if label not in _seen_labels and not _seen_labels.add(label)  # type: ignore[func-returns-value]
            ]

            agg_data: dict[int, dict[str, float]] = {g_idx: defaultdict(float) for g_idx, _ in group_cols}
            for row in probe:
                num_val: float | None = None
                if num_ci is not None and num_ci < len(row):
                    v = row[num_ci]
                    if isinstance(v, (int, float)):
                        num_val = float(v)
                if num_val is not None:
                    for g_idx, _ in group_cols:
                        gval = str(row[g_idx]) if g_idx < len(row) and row[g_idx] is not None else "(blank)"
                        agg_data[g_idx][gval] += num_val

            rev_col = headers[num_ci] if num_ci is not None and num_ci < len(headers) else "numeric"
            text_parts.append(f"Rows (sample): {len(probe)}  |  Revenue column: '{rev_col}'")
            for g_idx, g_label in group_cols:
                totals = dict(agg_data[g_idx])
                if not totals:
                    continue
                grand = sum(totals.values())
                text_parts.append(f"--- By {g_label} ---")
                for k, v in sorted(totals.items(), key=lambda x: -x[1])[:50]:
                    text_parts.append(f"  {k}: {v:,.2f}")
                text_parts.append(f"  TOTAL: {grand:,.2f}")
        else:
            text_parts.append(f"Columns: {', '.join(h for h in headers[:20] if h)}")
            for row in probe[:50]:
                text_parts.append(" | ".join(str(c) for c in row if c is not None))
        text_parts.append("")

        if len("\n".join(text_parts)) >= max_chars:
            break

    return "\n".join(text_parts)[:max_chars]


def _extract_sync(path: pathlib.Path, max_chars: int) -> str:
    """Synchronous dispatch — called from async via asyncio.to_thread()."""
    ext = path.suffix.lower()
    try:
        if ext in {".txt", ".md", ".csv"}:
            return _extract_txt(path, max_chars)
        elif ext == ".docx":
            return _extract_docx(path, max_chars)
        elif ext == ".pdf":
            return _extract_pdf(path, max_chars)
        elif ext in {".pptx", ".ppt"}:
            return _extract_pptx(path, max_chars)
        elif ext in {".xlsx", ".xls"}:
            return _extract_xlsx(path, max_chars)
        elif ext == ".boxnote":
            return _extract_boxnote(path, max_chars)
        elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
            return _extract_image(path, max_chars)
        else:
            return f"[Unsupported format: {path.name}]"
    except FileNotFoundError:
        logger.warning("File not found during extraction: %s", path)
        return f"[File not found: {path.name}]"
    except Exception as exc:
        logger.warning("Failed to extract %s: %s", path.name, exc)
        return f"[Extract error ({ext}): {exc}] File: {path.name}"


# ── Public async API ──────────────────────────────────────────────────────────

async def extract(
    file_path: str | pathlib.Path,
    max_chars: int | None = None,
) -> str:
    """
    Async extraction — returns up to *max_chars* characters of text.

    Defaults to KB_BUDGET_RAG_FILE chars (configured via KB_BUDGET_RAG_FILE env var).
    All file I/O runs in a thread pool so the event loop is never blocked.
    """
    path = pathlib.Path(file_path)
    limit = max_chars if max_chars is not None else _budget("rag_file")
    return await asyncio.to_thread(_extract_sync, path, limit)


def snippet(
    file_path: str | pathlib.Path,
    max_chars: int = 2000,
) -> str:
    """
    Synchronous short snippet extraction — used at index/generate time.
    Uses a smaller budget than full extraction.
    """
    path = pathlib.Path(file_path)
    return _extract_sync(path, max_chars)
