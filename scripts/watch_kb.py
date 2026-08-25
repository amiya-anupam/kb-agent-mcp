#!/usr/bin/env python3
"""
scripts/watch_kb.py — KnowledgeBase Dynamic Watcher
----------------------------------------------------
Watches KB_ROOT for filesystem events and keeps everything in sync:

  FILE ADDED      → embed into vector index + update README AUTO-INDEX
  FILE MODIFIED   → re-embed into vector index + update README AUTO-INDEX
  FILE DELETED    → remove from vector index + summary cache + update README
  FILE RENAMED    → remove old path from index/cache + embed new path + update README
  FOLDER CREATED  → run scripts/generate.py (builds index, creates agent, updates meta)
  FOLDER DELETED  → immediately delete agent .py + _index.json + meta entry
                    + all summary cache entries for that folder
  FOLDER RENAMED  → immediately rename agent .py + _index.json + meta entry
                    + re-key summary cache entries; no full scripts/generate.py needed

README AUTO-INDEX block (written into each folder's README):
  <!-- KB:AUTO-INDEX:START -->
  | file | type | size | modified | summary |
  ...per-file rows, always current...
  <!-- KB:AUTO-INDEX:END -->

Agents read this block directly — no raw file scanning at query time.

Configuration (env vars / .env):
  KB_ROOT           Root directory to watch (defaults to script directory)
  KB_MODEL          LLM model for generating file summaries
  KB_LLM_PROVIDER   ollama | openai | anthropic | custom
  KB_LLM_BASE_URL   LLM base URL
  KB_API_KEY        API key (openai / anthropic / custom)
  KB_IGNORE_FOLDERS Comma-separated extra folders to ignore

Run:
  python3 scripts/watch_kb.py
"""

import os
import re
import sys
import json
import time
import shutil
import subprocess
import datetime
import pathlib
import hashlib

from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# agents/ is in the repo root — resolve upward from scripts/ so agent_base
# and context_budget are importable at module level.
sys.path.insert(0, str(pathlib.Path(__file__).parent.parent / "agents"))
from context_budget import COLLAPSE_RULES, trim_summary, get as _budget_get
from agent_base import _find_readme as find_readme, DEFAULT_BLOCKLIST  # noqa: E402

sys.stdout.reconfigure(line_buffering=True)
sys.stderr.reconfigure(line_buffering=True)

# ── Environment loader ────────────────────────────────────────────────────────

def _load_env(root: pathlib.Path):
    env_file = root / ".env"
    if env_file.exists():
        for line in env_file.read_text(encoding="utf-8").splitlines():
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, _, v = line.partition("=")
                os.environ.setdefault(k.strip(), v.strip())

# scripts/ lives one level below the repo root — resolve upward so all
# relative paths (agents/, .env, etc.) stay correct.
SCRIPT_DIR = pathlib.Path(__file__).parent.parent.resolve()
_load_env(SCRIPT_DIR)

# ── Config (from env) ─────────────────────────────────────────────────────────

def _resolve_watch_root() -> pathlib.Path:
    raw = os.environ.get("KB_ROOT", "")
    return pathlib.Path(raw).resolve() if raw else SCRIPT_DIR

WATCH_ROOT   = _resolve_watch_root()
LLM_PROVIDER = os.environ.get("KB_LLM_PROVIDER", "ollama").lower()
LLM_BASE_URL = os.environ.get("KB_LLM_BASE_URL", "http://localhost:11434")
MODEL        = os.environ.get("KB_MODEL", "qwen3:14b")
API_KEY      = os.environ.get("KB_API_KEY", "")

# Stale-file threshold: warn when a file's mtime exceeds this many days.
# Set to 0 to disable staleness checks entirely.
_raw_stale = os.environ.get("KB_STALE_DAYS", "90")
try:
    STALE_DAYS = int(_raw_stale)
except ValueError:
    STALE_DAYS = 90

# How often (seconds) to re-run the staleness check while the watcher is live.
_STALE_CHECK_INTERVAL = 3600  # 1 hour

# How often (seconds) to force a full generate.py resync regardless of events.
# Guarantees eventual consistency even if individual file events were missed.
# Configurable via KB_RESYNC_HOURS (default 24 hours). Set to 0 to disable.
_raw_resync = os.environ.get("KB_RESYNC_HOURS", "24")
try:
    _RESYNC_INTERVAL = int(_raw_resync) * 3600
except ValueError:
    _RESYNC_INTERVAL = 24 * 3600

def _get_blocklist() -> set[str]:
    extra = os.environ.get("KB_IGNORE_FOLDERS", "")
    user  = {f.strip().lower() for f in extra.split(",") if f.strip()}
    return DEFAULT_BLOCKLIST | user

BLOCKLIST = _get_blocklist()

INCLUDE_EXTS  = {".pdf", ".docx", ".pptx", ".xlsx", ".ppt", ".doc",
                 ".png", ".jpg", ".jpeg", ".boxnote", ".md", ".txt", ".csv"}
# Files whose name contains any of these strings are never indexed or trigger README updates.
# "readme" catches all README variants; the watcher itself writes READMEs so we must
# ignore those writes to prevent an infinite update loop.
SKIP_PATTERNS = {"readme", ".ds_store", ".watch.log", "thumbs.db", "~$"}
DEBOUNCE_SECS = 5

# README markers
MARKER_START = "<!-- KB:AUTO-INDEX:START -->"
MARKER_END   = "<!-- KB:AUTO-INDEX:END -->"

# Summary cache path
SUMMARY_CACHE_PATH = SCRIPT_DIR / "agents" / "vector_store" / "file_summaries.json"
DOMAIN_META_PATH   = SCRIPT_DIR / "agents" / "vector_store" / "domain_meta.json"
AGENTS_DIR         = SCRIPT_DIR / "agents"
AUDIT_LOG_PATH     = SCRIPT_DIR / ".kb_index" / "audit.jsonl"

SUMMARY_EXTRACT_CHARS = 3000

# ── Folder name helpers (all dynamic — never hardcoded) ───────────────────────

def folder_to_safe_name(folder_name: str) -> str:
    """Convert any folder name to a safe snake_case identifier."""
    name = folder_name.lower()
    name = re.sub(r"[^a-z0-9]+", "_", name)
    return name.strip("_")


def agent_filename(folder_name: str) -> str:
    return f"agent_{folder_to_safe_name(folder_name)}.py"


def index_filename(folder_name: str) -> str:
    return f"{folder_to_safe_name(folder_name)}_index.json"


# ── Summary cache ─────────────────────────────────────────────────────────────

def _load_summary_cache() -> dict:
    if SUMMARY_CACHE_PATH.exists():
        try:
            return json.loads(SUMMARY_CACHE_PATH.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {}

def _save_summary_cache(cache: dict):
    SUMMARY_CACHE_PATH.parent.mkdir(parents=True, exist_ok=True)
    try:
        SUMMARY_CACHE_PATH.write_text(json.dumps(cache, indent=2), encoding="utf-8")
    except (PermissionError, OSError) as e:
        print(
            f"[KB Watcher] ✗ Could not write summary cache: {e}\n"
            f"  Path: {SUMMARY_CACHE_PATH}\n"
            f"  Summaries will be regenerated next run.",
            flush=True,
        )

# ── Domain meta helpers ───────────────────────────────────────────────────────

def _load_domain_meta() -> dict:
    if DOMAIN_META_PATH.exists():
        try:
            return json.loads(DOMAIN_META_PATH.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {}

def _save_domain_meta(meta: dict):
    DOMAIN_META_PATH.parent.mkdir(parents=True, exist_ok=True)
    try:
        DOMAIN_META_PATH.write_text(json.dumps(meta, indent=2, ensure_ascii=False), encoding="utf-8")
    except (PermissionError, OSError) as e:
        print(
            f"[KB Watcher] ✗ Could not write domain_meta.json: {e}\n"
            f"  Path: {DOMAIN_META_PATH}\n"
            f"  Domain routing changes will not persist until this is fixed.",
            flush=True,
        )

# ── Folder / file helpers ─────────────────────────────────────────────────────

def is_knowledge_folder(path: pathlib.Path) -> bool:
    return (
        path.is_dir()
        and path.parent == WATCH_ROOT
        and path.name.lower() not in BLOCKLIST
    )

def discover_knowledge_folders() -> list[pathlib.Path]:
    return [p for p in sorted(WATCH_ROOT.iterdir()) if is_knowledge_folder(p)]

# ── README finder — imported from agents/agent_base.py ───────────────────────
# find_readme() is the canonical copy in agent_base._find_readme().
# Imported at module level above (aliased to find_readme to preserve call sites).

def ensure_readme(folder: pathlib.Path) -> pathlib.Path:
    """
    Return the folder's README path, creating a minimal one if absent.
    Name is always derived from the folder name — never hardcoded.
    """
    existing = find_readme(folder)
    if existing:
        return existing
    readme = folder / f"{folder.name}.md"
    readme.write_text(
        f"# {folder.name}\n\n"
        f"Knowledge domain: **{folder.name}**\n\n"
        f"_Add your own notes about this domain here._\n\n",
        encoding="utf-8",
    )
    print(f"[KB Watcher] Created README: {readme.name}", flush=True)
    return readme

def _has_noindex_ancestor(path: pathlib.Path) -> bool:
    """
    Return True if any ancestor directory of *path* (up to WATCH_ROOT) contains
    a `.noindex` sentinel file.

    Mirrors the canonical implementation in kb_agent_mcp/file_parser and the
    agents-layer so all three layers enforce the same exclusion rule.
    """
    try:
        for parent in path.parents:
            if (parent / ".noindex").exists():
                return True
            if parent == WATCH_ROOT:
                break
    except Exception:
        pass
    return False


def should_skip(path: pathlib.Path) -> bool:
    """Return True if this file should be excluded from watcher processing.

    Skips:
    • files matching SKIP_PATTERNS (readme, .ds_store, etc.)
    • any file whose ancestor folder contains a `.noindex` sentinel file
    """
    if any(p in path.name.lower() for p in SKIP_PATTERNS):
        return True
    return _has_noindex_ancestor(path)

def is_readme(path: pathlib.Path) -> bool:
    """Return True if this path is a README file (should not be indexed as a doc)."""
    return "readme" in path.name.lower() or path.name.lower() == path.parent.name.lower() + ".md"

def human_size(size_bytes: int) -> str:
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 ** 2:
        return f"{size_bytes / 1024:.1f} KB"
    return f"{size_bytes / 1024 ** 2:.1f} MB"

def gather_files(folder: pathlib.Path) -> list[pathlib.Path]:
    return sorted(
        f for f in folder.rglob("*")
        if f.is_file()
        and f.suffix.lower() in INCLUDE_EXTS
        and not should_skip(f)
    )

def file_hash(path: pathlib.Path) -> str:
    try:
        return hashlib.md5(path.read_bytes()).hexdigest()
    except Exception:
        return ""

def top_folder_name(event_path: str) -> str | None:
    try:
        rel = pathlib.Path(event_path).relative_to(WATCH_ROOT)
        return rel.parts[0] if rel.parts else None
    except ValueError:
        return None

# ── Text extraction ───────────────────────────────────────────────────────────

def extract_snippet(file_path: pathlib.Path) -> str:
    """Extract up to SUMMARY_EXTRACT_CHARS of text from a file."""
    ext = file_path.suffix.lower()
    try:
        if ext in {".txt", ".md", ".csv"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")[:SUMMARY_EXTRACT_CHARS]

        elif ext == ".docx":
            import zipfile, re as _re
            with zipfile.ZipFile(file_path) as z:
                with z.open("word/document.xml") as f:
                    xml = f.read().decode("utf-8", errors="ignore")
            text = _re.sub(r"<[^>]+>", " ", xml)
            return " ".join(text.split())[:SUMMARY_EXTRACT_CHARS]

        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            text = ""
            for page in reader.pages[:4]:
                text += (page.extract_text() or "") + "\n"
                if len(text) >= SUMMARY_EXTRACT_CHARS:
                    break
            return text[:SUMMARY_EXTRACT_CHARS]

        elif ext in {".pptx", ".ppt"}:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides[:6]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text.strip() + "\n"
                if len(text) >= SUMMARY_EXTRACT_CHARS:
                    break
            return text[:SUMMARY_EXTRACT_CHARS]

        elif ext in {".xlsx", ".xls"}:
            # Large XLSX files (>50 MB) must use the streaming aggregator —
            # openpyxl takes 80+ seconds and produces only 30-row header dumps
            # which result in generic "XLSX file" summaries.  We delegate to
            # agent_base._stream_xlsx_aggregate() which produces a meaningful
            # revenue/data breakdown in 30-50s using raw XML iterparse.
            MAX_XLSX_BYTES = 50 * 1024 * 1024  # 50 MB
            if file_path.stat().st_size > MAX_XLSX_BYTES:
                try:
                    from agent_base import _stream_xlsx_aggregate
                    return _stream_xlsx_aggregate(file_path, SUMMARY_EXTRACT_CHARS)
                except Exception as e:
                    return f"[Large XLSX stream error: {e}]"

            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            text = ""
            for sheet in wb.worksheets[:2]:
                text += f"[Sheet: {sheet.title}] "
                for row in sheet.iter_rows(max_row=30, values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        text += row_text + " "
                if len(text) >= SUMMARY_EXTRACT_CHARS:
                    break
            return text[:SUMMARY_EXTRACT_CHARS]

        elif ext == ".boxnote":
            data = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
            def walk(node):
                if isinstance(node, dict):
                    if node.get("type") == "text":
                        yield node.get("text", "")
                    for v in node.values():
                        yield from walk(v)
                elif isinstance(node, list):
                    for item in node:
                        yield from walk(item)
            return " ".join(walk(data))[:SUMMARY_EXTRACT_CHARS]

    except Exception as e:
        return f"[extraction error: {e}]"

    return f"[{file_path.suffix.upper().lstrip('.')} file]"

# ── LLM helpers ───────────────────────────────────────────────────────────────

def _llm_available() -> bool:
    try:
        import httpx
        if LLM_PROVIDER == "ollama":
            r = httpx.get(f"{LLM_BASE_URL}/api/tags", timeout=4.0)
        else:
            r = httpx.get(LLM_BASE_URL.rstrip("/"), timeout=4.0)
        return r.status_code < 500
    except Exception:
        return False

def _call_llm(prompt: str) -> str:
    import httpx
    messages = [{"role": "user", "content": prompt}]

    if LLM_PROVIDER == "anthropic":
        headers = {"x-api-key": API_KEY, "anthropic-version": "2023-06-01",
                   "Content-Type": "application/json"}
        r = httpx.post(f"{LLM_BASE_URL}/v1/messages", headers=headers,
                       json={"model": MODEL, "max_tokens": 256,
                             "temperature": 0.1, "messages": messages},
                       timeout=30.0)
        r.raise_for_status()
        return r.json()["content"][0]["text"].strip()

    if LLM_PROVIDER in ("openai", "custom"):
        base = LLM_BASE_URL.rstrip("/")
        if "11434" in base and not base.endswith("/v1"):
            base += "/v1"
        headers = {"Content-Type": "application/json"}
        if API_KEY:
            headers["Authorization"] = f"Bearer {API_KEY}"
        r = httpx.post(f"{base}/chat/completions", headers=headers,
                       json={"model": MODEL, "messages": messages, "temperature": 0.1},
                       timeout=30.0)
        r.raise_for_status()
        return r.json()["choices"][0]["message"]["content"].strip()

    # Ollama
    r = httpx.post(f"{LLM_BASE_URL}/api/chat",
                   json={"model": MODEL, "messages": messages, "stream": False,
                         "options": {"temperature": 0.1, "num_ctx": _budget_get("num_ctx")},
                         "think": False},
                   timeout=30.0)
    r.raise_for_status()
    return r.json()["message"]["content"].strip()

# ── File summary (cached, LLM or heuristic fallback) ─────────────────────────
#
# Feature 8 — Hybrid heuristic summaries
# ──────────────────────────────────────
# When LLM is unavailable we generate a heuristic summary immediately (offline)
# and mark the cache entry with needs_llm_upgrade=True.  A background upgrade
# queue (max 5 per dispatch cycle) attempts the LLM call later; on success the
# cache entry is updated and the README re-written.

_LLM_UPGRADE_QUEUE: list[tuple[pathlib.Path, str]] = []   # (file_path, rel_key)
_LLM_UPGRADE_MAX_PER_CYCLE = 5


def _heuristic_summary(file_path: pathlib.Path, snippet: str) -> str:
    """
    Produce a best-effort offline summary without an LLM.

    Rules (in order):
    1. Large XLSX → keep the full streaming aggregate (numeric detail matters).
    2. Multi-line snippet → first meaningful line, stripped and capped.
    3. Short snippet → return as-is, capped to 200 chars.
    4. Nothing useful → "<EXT> file".
    """
    MAX_XLSX_BYTES = 50 * 1024 * 1024
    if (file_path.suffix.lower() in {".xlsx", ".xls"}
            and file_path.stat().st_size > MAX_XLSX_BYTES):
        return snippet  # preserve full structured summary

    lines = [l.strip() for l in snippet.splitlines() if l.strip()]
    if not lines:
        return f"{file_path.suffix.upper().lstrip('.')} file"

    # Skip generic slide/section markers like "[Slide 1]"
    first = next((l for l in lines if not re.match(r"^\[Slide \d+\]$", l)), lines[0])
    return trim_summary(first[:200], file_path.name)


def _llm_summary(file_path: pathlib.Path, snippet: str) -> str:
    """Generate a one-sentence LLM summary. Raises on failure."""
    prompt = (
        f"Summarise the following document excerpt in exactly ONE concise sentence "
        f"(max 20 words). Be specific about what the document covers.\n\n"
        f"Filename: {file_path.name}\n\nContent:\n{snippet[:1500]}\n\nOne-sentence summary:"
    )
    summary = _call_llm(prompt)
    summary = re.sub(
        r"^(summary|one.sentence summary|here is|this document)[:\s]+",
        "", summary, flags=re.IGNORECASE,
    ).strip()
    return trim_summary(summary, file_path.name)


def generate_file_summary(
    file_path: pathlib.Path,
    *,
    rel_key: str | None = None,
    upgrade_queue: list | None = None,
) -> tuple[str, bool]:
    """
    Return (summary, needs_llm_upgrade).

    When LLM is available: calls LLM, returns (summary, False).
    When LLM is unavailable: returns heuristic summary and (summary, True)
    so the caller can enqueue this file for a later LLM upgrade.

    *rel_key* and *upgrade_queue* are optional — when both are provided and
    an upgrade is needed, the (file_path, rel_key) pair is appended to
    upgrade_queue automatically.
    """
    snippet = extract_snippet(file_path)
    if not snippet or snippet.startswith("["):
        return f"{file_path.suffix.upper().lstrip('.')} file", False

    if _llm_available():
        try:
            return _llm_summary(file_path, snippet), False
        except Exception as e:
            print(f"[KB Watcher] ⚠ Summary LLM failed for {file_path.name}: {e}", flush=True)
            # Fall through to heuristic; mark for later upgrade
            summary = _heuristic_summary(file_path, snippet)
            if rel_key is not None and upgrade_queue is not None:
                upgrade_queue.append((file_path, rel_key))
            return summary, True

    # LLM unavailable — heuristic now, upgrade later
    summary = _heuristic_summary(file_path, snippet)
    if rel_key is not None and upgrade_queue is not None:
        upgrade_queue.append((file_path, rel_key))
    return summary, True


def get_file_summary(file_path: pathlib.Path, cache: dict, rel_key: str) -> tuple[str, bool]:
    """Return (summary, cache_updated). Reuses cached entry if hash unchanged."""
    h      = file_hash(file_path)
    cached = cache.get(rel_key, {})
    if cached.get("hash") == h and cached.get("summary"):
        # Re-queue for LLM upgrade if previously generated without one
        if cached.get("needs_llm_upgrade"):
            _LLM_UPGRADE_QUEUE.append((file_path, rel_key))
        return cached["summary"], False
    summary, needs_upgrade = generate_file_summary(
        file_path, rel_key=rel_key, upgrade_queue=_LLM_UPGRADE_QUEUE
    )
    cache[rel_key] = {"hash": h, "summary": summary, "needs_llm_upgrade": needs_upgrade}
    return summary, True

# ── AUTO-INDEX block builder ──────────────────────────────────────────────────

def build_auto_index_block(folder: pathlib.Path, cache: dict) -> tuple[str, bool]:
    """
    Build the raw AUTO-INDEX markdown block for a folder.

    Summary capping, collapse grouping, and all formatting rules are
    delegated to context_budget — this function only handles file I/O
    and cache management.
    """
    files       = gather_files(folder)
    now         = datetime.datetime.now().strftime("%d %b %Y %H:%M")
    cache_dirty = False
    summary_cap = _budget_get("summary")

    # ── Collect per-file (subdir, filename, summary) tuples ──────────────────
    collapsed: dict[str, dict] = {}   # label → {"files": [], "tmpl": str}
    row_data:  list[tuple]     = []   # (subdir, filename, summary)

    for f in files:
        rel     = f.relative_to(folder)
        subdir  = str(rel.parent) if len(rel.parts) > 1 else ""
        rel_key = str(f.relative_to(WATCH_ROOT))

        raw_summary, updated = get_file_summary(f, cache, rel_key)
        if updated:
            cache_dirty = True

        summary = trim_summary(raw_summary, f.name)

        # Check whether this file belongs to a collapse group
        matched = False
        for pat, label, tmpl in COLLAPSE_RULES:
            if pat.search(f.name):
                if label not in collapsed:
                    collapsed[label] = {"files": [], "tmpl": tmpl}
                collapsed[label]["files"].append(f.name)
                matched = True
                break

        if not matched:
            row_data.append((subdir, f.name, summary))

    # ── Build raw markdown table (compact_index_block will post-process) ─────
    lines = [
        "| File | Summary |",
        "|---|---|",
    ]

    current_subdir = None
    for subdir, fname, summary in row_data:
        if subdir != current_subdir:
            current_subdir = subdir
            if subdir:
                lines.append(f"| **📁 {subdir}/** | |")
        if len(summary) > summary_cap:
            summary = summary[:summary_cap] + "…"
        lines.append(f"| `{fname}` | {summary.replace('|', chr(92)+'|')} |")

    # One collapsed row per group
    for label, info in collapsed.items():
        n    = len(info["files"])
        qtrs = sorted(set(re.findall(r"Q[123]\d{2,3}", " ".join(info["files"]))))
        q_str = ", ".join(qtrs) if qtrs else f"{n} files"
        desc  = info["tmpl"].format(quarters=q_str, n=n)
        if len(desc) > summary_cap:
            desc = desc[:summary_cap] + "…"
        lines.append(f"| _{label}_ ({n} files) | {desc} |")

    lines.append("")
    return "\n".join(lines), cache_dirty


def update_readme(folder: pathlib.Path, cache: dict) -> bool:
    """
    Create or update the AUTO-INDEX block in the folder's README.
    Returns True if the summary cache was updated.
    """
    readme                 = ensure_readme(folder)
    content                = readme.read_text(encoding="utf-8")
    new_block, cache_dirty = build_auto_index_block(folder, cache)
    full_block             = MARKER_START + "\n\n" + new_block + "\n" + MARKER_END

    if MARKER_START in content and MARKER_END in content:
        pattern = re.escape(MARKER_START) + r".*?" + re.escape(MARKER_END)
        updated = re.sub(pattern, full_block, content, flags=re.DOTALL)
    else:
        updated = content.rstrip() + "\n\n---\n\n" + full_block + "\n"

    if updated != content:
        readme.write_text(updated, encoding="utf-8")
        print(
            f"[KB Watcher] ✓ README updated: {folder.name}/{readme.name} "
            f"({len(gather_files(folder))} files)",
            flush=True,
        )

    return cache_dirty

# ── Inline folder cleanup (no generate.py needed) ─────────────────────────────

def _purge_folder_artifacts(folder_name: str, cache: dict) -> bool:
    """
    Immediately remove all artifacts for a deleted/renamed folder:
      - agents/vector_store/<safe>_index.json  (vector embeddings)
      - domain_meta.json entry                 (routing + description)
      - all summary cache entries under that folder

    NOTE: There are no per-domain agent_*.py files in this architecture.
    The orchestrator is data-driven and reads domain_meta.json at runtime.
    Returns True if the summary cache was mutated.
    """
    safe        = folder_to_safe_name(folder_name)
    index_file  = AGENTS_DIR / "vector_store" / f"{safe}_index.json"
    cache_dirty = False

    if index_file.exists():
        index_file.unlink()
        print(f"[KB Watcher] 🗑 Deleted index: {index_file.name}", flush=True)

    # Remove domain_meta entry
    meta = _load_domain_meta()
    if folder_name in meta:
        del meta[folder_name]
        _save_domain_meta(meta)
        print(f"[KB Watcher] 🗑 Removed meta entry: {folder_name}", flush=True)

    # Purge all summary cache entries that belong to this folder
    prefix = folder_name + os.sep
    stale  = [k for k in cache if k.startswith(prefix) or k.startswith(folder_name + "/")]
    for k in stale:
        del cache[k]
    if stale:
        cache_dirty = True
        print(f"[KB Watcher] 🗑 Purged {len(stale)} cache entries for {folder_name}", flush=True)

    return cache_dirty


def _rename_folder_artifacts(old_name: str, new_name: str, cache: dict) -> bool:
    """
    Rename all artifacts when a top-level folder is renamed:
      - Rename <old>_index.json → <new>_index.json  (rewrite folder/path fields)
      - Update domain_meta.json key + folder_name / safe_name / agent_name fields
      - Re-key all summary cache entries

    NOTE: There are no per-domain agent_*.py files in this architecture.
    The orchestrator is data-driven and reads domain_meta.json at runtime.
    Returns True if the summary cache was mutated.
    """
    old_safe = folder_to_safe_name(old_name)
    new_safe = folder_to_safe_name(new_name)

    old_index = AGENTS_DIR / "vector_store" / f"{old_safe}_index.json"
    new_index = AGENTS_DIR / "vector_store" / f"{new_safe}_index.json"

    # Rename / update index file
    if old_index.exists():
        try:
            data = json.loads(old_index.read_text())
            # Update folder name and all rel_path keys
            data["folder"] = new_name
            for entry in data.get("entries", []):
                old_prefix = old_name + os.sep
                new_prefix = new_name + os.sep
                if entry.get("path", "").startswith(old_prefix):
                    entry["path"]   = new_prefix + entry["path"][len(old_prefix):]
                entry["folder"] = new_name
            new_index.write_text(json.dumps(data, indent=2), encoding="utf-8")
        except Exception as e:
            print(f"[KB Watcher] ⚠ Could not migrate index: {e}", flush=True)
        old_index.unlink()
        print(f"[KB Watcher] ✏ Renamed index: {old_index.name} → {new_index.name}", flush=True)

    # Update domain_meta.json
    meta = _load_domain_meta()
    if old_name in meta:
        entry        = meta.pop(old_name)
        entry["folder_name"] = new_name
        entry["safe_name"]   = new_safe
        entry["agent_name"]  = f"{new_name} Agent"
        meta[new_name]       = entry
        _save_domain_meta(meta)
        print(f"[KB Watcher] ✏ Updated meta: {old_name} → {new_name}", flush=True)

    # Re-key summary cache entries
    old_prefix = old_name + os.sep
    new_prefix = new_name + os.sep
    old_prefix_fwd = old_name + "/"
    new_prefix_fwd = new_name + "/"
    rekeyed    = {}
    cache_dirty = False
    for k, v in list(cache.items()):
        if k.startswith(old_prefix) or k.startswith(old_prefix_fwd):
            new_key         = new_prefix + k[len(old_prefix):] if k.startswith(old_prefix) \
                              else new_prefix_fwd + k[len(old_prefix_fwd):]
            rekeyed[new_key] = v
            del cache[k]
            cache_dirty = True
    cache.update(rekeyed)
    if cache_dirty:
        print(f"[KB Watcher] ✏ Re-keyed {len(rekeyed)} cache entries: {old_name} → {new_name}",
              flush=True)

    return cache_dirty

# ── Stale-file checker ────────────────────────────────────────────────────────

def check_stale_files(folders: list[pathlib.Path]) -> list[str]:
    """
    Scan all indexed files across the given knowledge folders and return a list
    of human-readable warning strings for any file whose mtime is older than
    STALE_DAYS days, e.g.:

        "⚠ BizOps: Q3_Renewal_Tracker.xlsx was last updated 112 days ago."

    Returns an empty list when STALE_DAYS is 0 (feature disabled) or when
    no files exceed the threshold.
    """
    if STALE_DAYS <= 0:
        return []

    now      = datetime.datetime.now()
    cutoff   = datetime.timedelta(days=STALE_DAYS)
    warnings: list[str] = []

    for folder in folders:
        folder_name = folder.name
        for f in gather_files(folder):
            try:
                mtime     = datetime.datetime.fromtimestamp(f.stat().st_mtime)
                age       = now - mtime
                age_days  = age.days
            except OSError:
                continue

            if age >= cutoff:
                warnings.append(
                    f"⚠ {folder_name}: {f.name} was last updated {age_days} days ago."
                )

    return warnings


# ── Audit log ─────────────────────────────────────────────────────────────────

def _write_audit(event: str, detail: dict | None = None):
    """Append a JSON entry to .kb_index/audit.jsonl."""
    try:
        AUDIT_LOG_PATH.parent.mkdir(parents=True, exist_ok=True)
        entry = {
            "ts":    datetime.datetime.now().isoformat(timespec="seconds"),
            "event": event,
        }
        if detail:
            entry.update(detail)
        with AUDIT_LOG_PATH.open("a", encoding="utf-8") as fh:
            fh.write(json.dumps(entry) + "\n")
    except Exception as e:
        print(f"[KB Watcher] ⚠ Could not write audit log: {e}", flush=True)


# ── generate.py trigger ───────────────────────────────────────────────────────

# Configurable via KB_GENERATE_TIMEOUT (seconds). Default 900s (15 min) to
# handle large folders with many files needing LLM summarisation. The old
# hard-coded 300s was too short for BizOps (72 files × ~4s each ≈ 288s minimum).
_raw_gen_timeout = os.environ.get("KB_GENERATE_TIMEOUT", "900")
try:
    _GENERATE_TIMEOUT = int(_raw_gen_timeout)
except ValueError:
    _GENERATE_TIMEOUT = 900


# ── Per-domain incremental reindex ────────────────────────────────────────────

def run_reindex(domain: str, reason: str):
    """
    Re-embed all changed/new files in a single domain without running the full
    generate.py pipeline.

    Uses agents/embeddings.build_index() directly (hash-skips unchanged files)
    so only files that actually changed are re-embedded.  A post-rescan then
    catches any file that arrived while the embedding loop was running.

    Appropriate for file add/modify/delete events inside an *existing* domain.
    New top-level folders still go through run_generate() so domain_meta.json
    and SKILL.md are updated.
    """
    print(
        f"[KB Watcher] 🔄 Reindexing domain: {domain} ({reason})",
        flush=True,
    )
    start_time = time.time()

    sys.path.insert(0, str(AGENTS_DIR))
    try:
        from embeddings import build_index
        result = build_index(domain)
        n = len(result.get("entries", []))
        print(f"[KB Watcher] ✓ Reindex complete: {domain} ({n} files)", flush=True)

        # Write audit entry so drift-check CLI can compare against folder state.
        folders = discover_knowledge_folders()
        _write_audit("reindex_complete", {
            "reason": reason,
            "domain": domain,
            "files":  n,
        })

        # Catch files that arrived during the embedding loop
        _rescan_domain_since(domain, start_time)
    except Exception as e:
        print(f"[KB Watcher] ✗ Reindex failed for {domain}: {e}", flush=True)


def _rescan_domain_since(domain: str, since: float):
    """
    Re-index any file in *domain* whose mtime is newer than *since*.
    Closes the race window between the start of build_index() and its completion.
    """
    folder = WATCH_ROOT / domain
    if not folder.exists():
        return
    caught = 0
    for f in gather_files(folder):
        try:
            if f.stat().st_mtime >= since:
                print(
                    f"[KB Watcher] 🔁 Post-reindex rescan: {f.name} in {domain}",
                    flush=True,
                )
                _update_index(domain, f)
                caught += 1
        except OSError:
            continue
    if caught:
        print(f"[KB Watcher] ✓ Post-reindex rescan caught {caught} file(s) in {domain}", flush=True)


def run_generate(reason: str):
    generate_script = SCRIPT_DIR / "scripts" / "generate.py"
    if not generate_script.exists():
        print(f"[KB Watcher] ⚠ scripts/generate.py not found — skipping agent rebuild", flush=True)
        return

    print(
        f"[KB Watcher] 🔄 Running scripts/generate.py ({reason}, "
        f"timeout={_GENERATE_TIMEOUT}s)...",
        flush=True,
    )

    # Record start time so we can catch files that arrived during the generate window.
    start_time = time.time()

    try:
        # Forward KB_ROOT explicitly so the subprocess always knows which root
        # to scan — even if the parent process set it via a .env file that the
        # child's setdefault() would not override from an inherited env.
        env = os.environ.copy()
        env["KB_ROOT"] = str(WATCH_ROOT)
        result = subprocess.run(
            [sys.executable, str(generate_script)],
            capture_output=False,
            cwd=str(SCRIPT_DIR),
            env=env,
            timeout=_GENERATE_TIMEOUT,
        )
        if result.returncode == 0:
            print(f"[KB Watcher] ✓ generate.py completed", flush=True)
            # Write audit entry so drift-check CLI can compare against folder state.
            folders = discover_knowledge_folders()
            total_files = sum(len(gather_files(f)) for f in folders)
            _write_audit("generate_complete", {
                "reason":  reason,
                "domains": [f.name for f in folders],
                "files":   total_files,
            })
            # Re-scan for any files that arrived or changed during the generate window.
            # This closes the race condition where a user drops files while generate.py
            # is running and those files are not in generate.py's snapshot.
            _rescan_since(start_time)
        else:
            print(f"[KB Watcher] ✗ generate.py exited {result.returncode}", flush=True)
    except subprocess.TimeoutExpired:
        print(
            f"[KB Watcher] ✗ generate.py timed out after {_GENERATE_TIMEOUT}s "
            f"(set KB_GENERATE_TIMEOUT env var to increase)",
            flush=True,
        )
    except Exception as e:
        print(f"[KB Watcher] ✗ generate.py failed: {e}", flush=True)


def _rescan_since(since: float):
    """
    Walk all knowledge folders and re-index any file whose mtime is newer than
    `since` (a Unix timestamp captured just before generate.py was invoked).

    This closes the race window: files added or replaced while generate.py was
    running are caught here and embedded into the vector store immediately.
    """
    caught = 0
    for folder in discover_knowledge_folders():
        folder_name = folder.name
        for f in gather_files(folder):
            try:
                if f.stat().st_mtime >= since:
                    print(
                        f"[KB Watcher] 🔁 Post-generate rescan: {f.name} in {folder_name}",
                        flush=True,
                    )
                    _update_index(folder_name, f)
                    caught += 1
            except OSError:
                continue
    if caught:
        print(f"[KB Watcher] ✓ Post-generate rescan caught {caught} file(s)", flush=True)
    else:
        print(f"[KB Watcher] ✓ Post-generate rescan: nothing missed", flush=True)

# ── Embeddings helper (lazy import) ───────────────────────────────────────────

def _update_index(folder_name: str, file_path: pathlib.Path):
    """Embed a single file into its folder's vector index."""
    sys.path.insert(0, str(AGENTS_DIR))
    try:
        from embeddings import update_index_for_file
        updated = update_index_for_file(folder_name, file_path)
        if updated:
            print(f"[KB Watcher] 🔍 Indexed: {file_path.name}", flush=True)
    except Exception as e:
        print(f"[KB Watcher] ⚠ Index update failed for {file_path.name}: {e}", flush=True)


def _remove_index(folder_name: str, file_path: pathlib.Path):
    """Remove a single file from its folder's vector index."""
    sys.path.insert(0, str(AGENTS_DIR))
    try:
        from embeddings import remove_from_index
        removed = remove_from_index(folder_name, file_path)
        if removed:
            print(f"[KB Watcher] 🔍 Removed from index: {file_path.name}", flush=True)
    except Exception as e:
        print(f"[KB Watcher] ⚠ Index removal failed for {file_path.name}: {e}", flush=True)

# ── Event handler ─────────────────────────────────────────────────────────────

class KBHandler(FileSystemEventHandler):
    """
    Pending work queues (debounced):
      _pending_readme[folder_name] = deadline   → update README + re-index files in that folder
      _pending_generate = deadline              → run generate.py (new folder only)

    Immediate work (no debounce):
      folder deleted → _purge_folder_artifacts()
      folder renamed → _rename_folder_artifacts()
    """

    def __init__(self):
        super().__init__()
        self._pending_readme:           dict[str, float] = {}
        self._pending_index:            dict[str, set[pathlib.Path]] = {}  # folder → files to index
        self._pending_deindex:          dict[str, set[pathlib.Path]] = {}  # folder → files to remove
        self._pending_generate:         float | None = None
        self._pending_generate_reason:  str          = ""
        self._pending_reindex:          dict[str, str] = {}   # domain → reason (debounced per-domain)
        self._known_folders:            set[str]     = {
            p.name for p in discover_knowledge_folders()
        }
        self._cache:                    dict         = _load_summary_cache()
        self._next_stale_check:         float        = time.time() + _STALE_CHECK_INTERVAL
        self._next_resync:              float        = (
            time.time() + _RESYNC_INTERVAL if _RESYNC_INTERVAL > 0 else float("inf")
        )

    # ── Schedulers ────────────────────────────────────────────────────────────

    def _schedule_readme(self, folder_name: str):
        self._pending_readme[folder_name] = time.time() + DEBOUNCE_SECS

    def _schedule_index(self, folder_name: str, file_path: pathlib.Path):
        self._pending_index.setdefault(folder_name, set()).add(file_path)
        self._schedule_readme(folder_name)

    def _schedule_deindex(self, folder_name: str, file_path: pathlib.Path):
        self._pending_deindex.setdefault(folder_name, set()).add(file_path)
        self._schedule_readme(folder_name)

    def _schedule_generate(self, reason: str):
        if self._pending_generate is None:
            self._pending_generate        = time.time() + DEBOUNCE_SECS
            self._pending_generate_reason = reason

    def _schedule_reindex(self, domain: str, reason: str):
        """
        Schedule a per-domain reindex after the debounce window.

        Multiple events for the same domain within DEBOUNCE_SECS are coalesced
        into a single build_index() call — the last reason wins (descriptive only).
        The debounce is shared with the README update: the reindex fires after
        the same 10 s window so embeddings and README stay in sync.
        """
        self._pending_reindex[domain] = reason
        # Re-use the README debounce deadline so both fire together
        self._schedule_readme(domain)

    # ── Dispatch loop (called every second) ──────────────────────────────────

    def dispatch_pending(self):
        now = time.time()

        # 1. Process deindex requests first (so removed files don't appear in README)
        for folder_name, files in list(self._pending_deindex.items()):
            # Only fire when the matching readme debounce has passed
            if self._pending_readme.get(folder_name, 0) <= now:
                del self._pending_deindex[folder_name]
                for fp in files:
                    _remove_index(folder_name, fp)
                    rel_key = str(fp.relative_to(WATCH_ROOT))
                    if rel_key in self._cache:
                        del self._cache[rel_key]
                _save_summary_cache(self._cache)

        # 2. Process index requests
        for folder_name, files in list(self._pending_index.items()):
            if self._pending_readme.get(folder_name, 0) <= now:
                del self._pending_index[folder_name]
                for fp in files:
                    if fp.exists():
                        _update_index(folder_name, fp)

        # 3. Fire overdue README updates + per-domain reindexes together
        fired = [fn for fn, t in self._pending_readme.items() if now >= t]
        for folder_name in fired:
            del self._pending_readme[folder_name]
            folder = WATCH_ROOT / folder_name
            if folder.exists() and is_knowledge_folder(folder):
                dirty = update_readme(folder, self._cache)
                if dirty:
                    _save_summary_cache(self._cache)

            # Fire any pending per-domain reindex for this folder
            if folder_name in self._pending_reindex:
                reason = self._pending_reindex.pop(folder_name)
                run_reindex(folder_name, reason)

        # 4. Poll for folder changes the OS may not have delivered directory events for.
        #    macOS FSEvents sometimes coalesces mkdir/rmdir into just file events.
        live_folders    = {p.name for p in discover_knowledge_folders()}
        new_folders     = live_folders - self._known_folders
        deleted_folders = self._known_folders - live_folders

        for fname in new_folders:
            self._known_folders.add(fname)
            print(f"[KB Watcher] 📁 New folder detected (poll): {fname} → scheduling generate.py",
                  flush=True)
            self._schedule_generate(f"new folder: {fname}")

        for fname in deleted_folders:
            self._known_folders.discard(fname)
            print(f"[KB Watcher] 🗑 Folder gone (poll): {fname} → purging artifacts immediately",
                  flush=True)
            dirty = _purge_folder_artifacts(fname, self._cache)
            if dirty:
                _save_summary_cache(self._cache)
            # Cancel any pending work for this folder
            self._pending_readme.pop(fname, None)
            self._pending_index.pop(fname, None)
            self._pending_deindex.pop(fname, None)
            self._pending_reindex.pop(fname, None)

        # 5. Fire overdue generate.py run (new folder only)
        if self._pending_generate is not None and now >= self._pending_generate:
            reason                 = self._pending_generate_reason
            self._pending_generate = None
            run_generate(reason)
            self._known_folders = {p.name for p in discover_knowledge_folders()}

        # 6. Periodic force-resync (default every 24 hours)
        #    Guarantees eventual consistency even if individual file events were dropped.
        if _RESYNC_INTERVAL > 0 and now >= self._next_resync:
            self._next_resync = now + _RESYNC_INTERVAL
            print(f"[KB Watcher] 🔄 Scheduled 24h resync — running generate.py", flush=True)
            run_generate("scheduled resync")
            self._known_folders = {p.name for p in discover_knowledge_folders()}

        # 7. Hourly stale-file check
        if STALE_DAYS > 0 and now >= self._next_stale_check:
            self._next_stale_check = now + _STALE_CHECK_INTERVAL
            folders = [WATCH_ROOT / fn for fn in self._known_folders
                       if (WATCH_ROOT / fn).is_dir()]
            stale   = check_stale_files(folders)
            if stale:
                print("[KB Watcher] 🕐 Stale file check:", flush=True)
                for w in stale:
                    print(f"  {w}", flush=True)

        # 8. LLM upgrade queue — attempt to upgrade heuristic summaries to LLM
        #    quality when the LLM becomes reachable.  Process at most
        #    _LLM_UPGRADE_MAX_PER_CYCLE entries per dispatch cycle to avoid
        #    blocking the event loop for too long.
        if _LLM_UPGRADE_QUEUE and _llm_available():
            upgraded_folders: set[str] = set()
            batch = _LLM_UPGRADE_QUEUE[:_LLM_UPGRADE_MAX_PER_CYCLE]
            del _LLM_UPGRADE_QUEUE[:_LLM_UPGRADE_MAX_PER_CYCLE]
            for fp, rel_key in batch:
                if not fp.exists():
                    continue
                cached = self._cache.get(rel_key, {})
                if not cached.get("needs_llm_upgrade"):
                    continue  # already upgraded by a prior cycle
                snippet = extract_snippet(fp)
                if not snippet or snippet.startswith("["):
                    continue
                try:
                    summary = _llm_summary(fp, snippet)
                    self._cache[rel_key] = {
                        "hash": cached.get("hash", ""),
                        "summary": summary,
                        "needs_llm_upgrade": False,
                    }
                    folder_name = fp.relative_to(WATCH_ROOT).parts[0]
                    upgraded_folders.add(folder_name)
                    print(f"[KB Watcher] ✨ LLM upgrade: {fp.name}", flush=True)
                except Exception as e:
                    print(f"[KB Watcher] ⚠ LLM upgrade failed for {fp.name}: {e}", flush=True)
                    _LLM_UPGRADE_QUEUE.append((fp, rel_key))  # re-queue for next cycle
            if upgraded_folders:
                _save_summary_cache(self._cache)
                for folder_name in upgraded_folders:
                    folder = WATCH_ROOT / folder_name
                    if folder.exists() and is_knowledge_folder(folder):
                        update_readme(folder, self._cache)
                        print(f"[KB Watcher] ✓ README refreshed after LLM upgrades: {folder_name}",
                              flush=True)

    # ── watchdog callbacks ────────────────────────────────────────────────────

    def on_created(self, event):
        path = pathlib.Path(event.src_path)

        if event.is_directory:
            if path.parent == WATCH_ROOT and path.name.lower() not in BLOCKLIST:
                if path.name not in self._known_folders:
                    self._known_folders.add(path.name)
                    print(f"[KB Watcher] 📁 New folder: {path.name} → scheduling generate.py",
                          flush=True)
                    self._schedule_generate(f"new folder: {path.name}")
            return

        folder_name = top_folder_name(event.src_path)
        if not folder_name or folder_name.lower() in BLOCKLIST:
            return
        if is_readme(path) or should_skip(path):
            return
        if path.suffix.lower() not in INCLUDE_EXTS:
            return

        print(f"[KB Watcher] ＋ {path.name} in {folder_name} → scheduling index + README",
              flush=True)
        self._schedule_index(folder_name, path)

        # File created inside a sub-folder — schedule a per-domain reindex so
        # the new file's embedding is added without a full generate.py run.
        if path.parent != WATCH_ROOT / folder_name:
            self._schedule_reindex(folder_name, f"sub-folder content added: {folder_name}")

    def on_deleted(self, event):
        path = pathlib.Path(event.src_path)

        if event.is_directory:
            if path.parent == WATCH_ROOT and path.name in self._known_folders:
                self._known_folders.discard(path.name)
                print(f"[KB Watcher] 🗑 Folder deleted: {path.name} → purging artifacts immediately",
                      flush=True)
                dirty = _purge_folder_artifacts(path.name, self._cache)
                if dirty:
                    _save_summary_cache(self._cache)
                # Cancel any pending work for this folder
                self._pending_readme.pop(path.name, None)
                self._pending_index.pop(path.name, None)
                self._pending_deindex.pop(path.name, None)
                self._pending_reindex.pop(path.name, None)
            return

        folder_name = top_folder_name(event.src_path)
        if not folder_name or folder_name.lower() in BLOCKLIST:
            return
        if is_readme(path) or path.suffix.lower() not in INCLUDE_EXTS:
            return

        print(f"[KB Watcher] ✕ {path.name} deleted from {folder_name} → scheduling deindex + README",
              flush=True)
        self._schedule_deindex(folder_name, path)

        # File deleted from a sub-folder — schedule a per-domain reindex so
        # the removed file is pruned from the index without a full generate.py run.
        if path.parent != WATCH_ROOT / folder_name:
            self._schedule_reindex(folder_name, f"sub-folder content removed: {folder_name}")

    def on_moved(self, event):
        src  = pathlib.Path(event.src_path)
        dest = pathlib.Path(event.dest_path)

        # ── Top-level folder renamed ──────────────────────────────────────────
        if event.is_directory and src.parent == WATCH_ROOT and dest.parent == WATCH_ROOT:
            old_name = src.name
            new_name = dest.name

            if old_name in self._known_folders and new_name.lower() not in BLOCKLIST:
                self._known_folders.discard(old_name)
                self._known_folders.add(new_name)
                print(f"[KB Watcher] ✏ Folder renamed: {old_name} → {new_name} "
                      f"→ updating artifacts immediately", flush=True)

                dirty = _rename_folder_artifacts(old_name, new_name, self._cache)
                if dirty:
                    _save_summary_cache(self._cache)

                # Update README (the renamed folder now exists at dest)
                self._schedule_readme(new_name)

                # Migrate pending work from old name to new name
                for q in (self._pending_readme, self._pending_index,
                          self._pending_deindex, self._pending_reindex):
                    if old_name in q:
                        q[new_name] = q.pop(old_name)
            return

        # ── File moved / renamed inside a knowledge folder ────────────────────
        src_folder  = top_folder_name(event.src_path)
        dest_folder = top_folder_name(event.dest_path)

        if src_folder and src_folder.lower() not in BLOCKLIST:
            if not is_readme(src) and src.suffix.lower() in INCLUDE_EXTS and not should_skip(src):
                print(f"[KB Watcher] ↩ {src.name} moved/renamed → deindexing src", flush=True)
                self._schedule_deindex(src_folder, src)

        if dest_folder and dest_folder.lower() not in BLOCKLIST:
            if not is_readme(dest) and dest.suffix.lower() in INCLUDE_EXTS and not should_skip(dest):
                print(f"[KB Watcher] ↪ {dest.name} moved/renamed → indexing dest", flush=True)
                self._schedule_index(dest_folder, dest)

    def on_modified(self, event):
        if event.is_directory:
            return
        path = pathlib.Path(event.src_path)
        # Ignore README writes — we write READMEs ourselves; don't loop
        if is_readme(path):
            return
        folder_name = top_folder_name(event.src_path)
        if not folder_name or folder_name.lower() in BLOCKLIST:
            return
        if path.suffix.lower() not in INCLUDE_EXTS or should_skip(path):
            return

        print(f"[KB Watcher] ✎ {path.name} modified in {folder_name} → scheduling re-index + README",
              flush=True)
        self._schedule_index(folder_name, path)

        # File modified inside a sub-folder — schedule a per-domain reindex so
        # the updated embedding is stored without a full generate.py run.
        if path.parent != WATCH_ROOT / folder_name:
            self._schedule_reindex(folder_name, f"sub-folder content changed: {folder_name}")

# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    if not WATCH_ROOT.exists():
        env_file = SCRIPT_DIR / ".env"
        print(f"[KB Watcher] ✗ WATCH_ROOT does not exist: {WATCH_ROOT}", flush=True)
        if env_file.exists():
            print(f"[KB Watcher]   Check KB_ROOT in your .env: {env_file}", flush=True)
        else:
            print(f"[KB Watcher]   Set KB_ROOT in {SCRIPT_DIR}/.env", flush=True)
        print(f"[KB Watcher]   Or run:  python3 scripts/setup.py   to configure interactively.", flush=True)
        sys.exit(1)

    print(f"[KB Watcher] Starting — watching {WATCH_ROOT}", flush=True)
    print(f"[KB Watcher] LLM: {LLM_PROVIDER} / {MODEL}", flush=True)
    if not _llm_available():
        print(f"[KB Watcher] ⚠ LLM not reachable — summaries will use text snippets as fallback",
              flush=True)

    folders = discover_knowledge_folders()
    if folders:
        for folder in folders:
            print(f"  → {folder.name}/ ({len(gather_files(folder))} files)", flush=True)
    else:
        print("  ⚠ No knowledge folders found — will watch for new ones", flush=True)

    # Startup stale-file check
    if STALE_DAYS > 0:
        print(
            f"[KB Watcher] Stale threshold: {STALE_DAYS} days "
            f"(KB_STALE_DAYS — set to 0 to disable)",
            flush=True,
        )
        stale_warnings = check_stale_files(folders)
        if stale_warnings:
            print(f"[KB Watcher] ⚠ Stale files detected:", flush=True)
            for w in stale_warnings:
                print(f"  {w}", flush=True)
        else:
            print(f"[KB Watcher] ✓ No stale files (all within {STALE_DAYS} days)", flush=True)

    handler  = KBHandler()
    observer = Observer()
    observer.schedule(handler, str(WATCH_ROOT), recursive=True)
    observer.start()
    print("[KB Watcher] Running. Press Ctrl+C to stop.\n", flush=True)

    try:
        while True:
            handler.dispatch_pending()
            time.sleep(1)
    except KeyboardInterrupt:
        print("\n[KB Watcher] Stopping…", flush=True)
        observer.stop()
    observer.join()
    print("[KB Watcher] Stopped.", flush=True)


if __name__ == "__main__":
    main()
