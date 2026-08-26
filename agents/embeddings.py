#!/usr/bin/env python3
"""
embeddings.py — Dynamic vector index builder and semantic search
----------------------------------------------------------------
Supports three embedding backends (tried in order):
  1. Ollama (KB_EMBED_MODEL via KB_LLM_BASE_URL)
  2. OpenAI-compatible API (when KB_LLM_PROVIDER=openai|custom)
  3. sentence-transformers offline fallback (all-MiniLM-L6-v2)

Folders are auto-discovered from KB_ROOT — no hardcoding required.
Index files are stored in agents/vector_store/<safe_name>_index.json.

Run standalone to rebuild all indexes:
  python3 agents/embeddings.py [--force]
"""

import os
import sys
import json
import pathlib
import hashlib
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity

# ── Config (all from env) ─────────────────────────────────────────────────────

def _kb_root() -> pathlib.Path:
    raw = os.environ.get("KB_ROOT", "")
    if raw:
        return pathlib.Path(raw)
    # Default: two levels up from this file (agents/embeddings.py → KB root)
    return pathlib.Path(__file__).parent.parent

# !! DO NOT REMOVE KB_ROOT from this module !!
# test_noindex_guard.py patches it via monkeypatch.setattr(embeddings, "KB_ROOT", tmp_path).
# _has_noindex_ancestor is imported from agent_base and reads agent_base.KB_ROOT internally,
# so the patch doesn't affect the boundary check — but removing the name causes AttributeError
# in those tests.
KB_ROOT      = _kb_root()
VECTOR_STORE = pathlib.Path(__file__).parent / "vector_store"

# Import shared helpers and XLSX config from agent_base (same agents/ directory)
sys.path.insert(0, str(pathlib.Path(__file__).parent))
import importlib as _importlib
_cb = _importlib.import_module("context_budget")
from agent_base import (
    _has_noindex_ancestor,
    folder_to_safe_name,
    should_skip,
    AGG_KEYWORDS,
    PREFERRED_NUM_COLS,
    DEFAULT_BLOCKLIST,
    INCLUDE_EXTS,
    SKIP_PATTERNS,
)

OLLAMA_URL   = os.environ.get("KB_LLM_BASE_URL", "http://localhost:11434")
EMBED_MODEL  = os.environ.get("KB_EMBED_MODEL", "nomic-embed-text")
LLM_PROVIDER = os.environ.get("KB_LLM_PROVIDER", "ollama").lower()
API_KEY      = os.environ.get("KB_API_KEY", "")

def _user_blocklist() -> set[str]:
    raw = os.environ.get("KB_IGNORE_FOLDERS", "")
    return {f.strip().lower() for f in raw.split(",") if f.strip()}

BLOCKLIST = DEFAULT_BLOCKLIST | _user_blocklist()

# Chars of text used as embedding input per file
SUMMARY_CHARS = 2000


def index_path_for(folder_name: str) -> pathlib.Path:
    safe = folder_to_safe_name(folder_name)
    return VECTOR_STORE / f"{safe}_index.json"


# ── Domain discovery ──────────────────────────────────────────────────────────

def discover_folders() -> list[str]:
    """
    Return all top-level subdirectory names under KB_ROOT that are not
    in the blocklist and contain at least one indexable file.
    """
    folders = []
    for p in sorted(KB_ROOT.iterdir()):
        if not p.is_dir():
            continue
        if p.name.lower() in BLOCKLIST:
            continue
        # Must contain at least one indexable file
        has_files = any(
            f.suffix.lower() in INCLUDE_EXTS
            for f in p.rglob("*")
            if f.is_file() and not should_skip(f)
        )
        if has_files:
            folders.append(p.name)
    return folders


# ── Embedding backends ────────────────────────────────────────────────────────

def _embed_ollama(text: str) -> list[float]:
    import httpx
    response = httpx.post(
        f"{OLLAMA_URL}/api/embeddings",
        json={"model": EMBED_MODEL, "prompt": text[:_cb.get("embed_chars")]},
        timeout=30.0,
    )
    response.raise_for_status()
    return response.json()["embedding"]


def _embed_openai(text: str) -> list[float]:
    import httpx
    base = os.environ.get("KB_LLM_BASE_URL", "https://api.openai.com/v1").rstrip("/")
    headers = {"Content-Type": "application/json"}
    if API_KEY:
        headers["Authorization"] = f"Bearer {API_KEY}"
    response = httpx.post(
        f"{base}/embeddings",
        headers=headers,
        json={"model": EMBED_MODEL or "text-embedding-3-small", "input": text[:_cb.get("embed_chars")]},
        timeout=30.0,
    )
    response.raise_for_status()
    return response.json()["data"][0]["embedding"]


_st_model = None  # lazy-loaded sentence-transformers model
_ST_MODEL_NAME = "all-MiniLM-L6-v2"

def _embed_sentence_transformers(text: str) -> list[float]:
    global _st_model
    if _st_model is None:
        try:
            from sentence_transformers import SentenceTransformer
            try:
                _st_model = SentenceTransformer(_ST_MODEL_NAME)
            except OSError as exc:
                # Raised when the model is not cached and the network is unreachable.
                raise RuntimeError(
                    f"Embedding model '{_ST_MODEL_NAME}' could not be loaded — "
                    "no internet connection and the model is not cached on this machine.\n"
                    "Run this once on a machine with internet access to pre-cache it:\n"
                    f"  python -c \"from sentence_transformers import SentenceTransformer; "
                    f"SentenceTransformer('{_ST_MODEL_NAME}')\"\n"
                    "After that the model loads entirely offline."
                ) from exc
        except ImportError:
            raise ImportError(
                "sentence-transformers is not installed and the configured embed "
                "endpoint is not reachable.\n"
                "Install it with:  pip install sentence-transformers"
            )
    return _st_model.encode(text[:_cb.get("embed_chars")]).tolist()


def get_embedding(text: str) -> list[float]:
    """
    Try embedding backends in order:
      1. passthrough mode → sentence-transformers directly (no LLM available)
      2. Configured provider (Ollama or OpenAI-compatible)
      3. sentence-transformers offline fallback
    """
    # passthrough mode has no LLM — go straight to sentence-transformers
    if LLM_PROVIDER == "passthrough":
        try:
            return _embed_sentence_transformers(text)
        except ImportError:
            raise RuntimeError(
                "KB_LLM_PROVIDER=passthrough requires sentence-transformers for embeddings.\n"
                "Install it with:  pip install sentence-transformers\n"
                "(This downloads ~80 MB the first time, then works fully offline.)"
            )

    # Try primary provider
    try:
        if LLM_PROVIDER in ("openai", "anthropic", "custom"):
            return _embed_openai(text)
        else:
            return _embed_ollama(text)
    except Exception as primary_err:
        # Fallback to sentence-transformers
        try:
            vec = _embed_sentence_transformers(text)
            return vec
        except ImportError:
            raise RuntimeError(
                f"Primary embedding failed ({primary_err}) and sentence-transformers "
                f"is not installed.\nEither fix the LLM connection or run:\n"
                f"  pip install sentence-transformers"
            )


# ── Text extraction (lightweight — for embedding summaries only) ──────────────

def extract_text_snippet(file_path: pathlib.Path) -> str:
    """Extract a short text snippet from a file for embedding."""
    ext = file_path.suffix.lower()
    try:
        if ext in {".txt", ".md", ".csv"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")[:SUMMARY_CHARS]

        elif ext == ".docx":
            import zipfile, re
            with zipfile.ZipFile(file_path) as z:
                with z.open("word/document.xml") as f:
                    xml = f.read().decode("utf-8", errors="ignore")
            text = re.sub(r"<[^>]+>", " ", xml)
            return " ".join(text.split())[:SUMMARY_CHARS]

        elif ext == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(str(file_path))
                text = ""
                for page in reader.pages[:3]:
                    text += page.extract_text() or ""
                    if len(text) >= SUMMARY_CHARS:
                        break
                return text[:SUMMARY_CHARS]
            except Exception:
                return f"PDF: {file_path.name}"

        elif ext in {".pptx", ".ppt"}:
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                text = ""
                for slide in prs.slides[:5]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
                    if len(text) >= SUMMARY_CHARS:
                        break
                return text[:SUMMARY_CHARS]
            except Exception:
                return f"PPTX: {file_path.name}"

        elif ext in {".xlsx", ".xls"}:
            try:
                import openpyxl

                # ── Skip very large XLSX files (>50 MB) ───────────────────────
                # openpyxl needs to decompress the entire ZIP — a 160 MB file
                # takes 80+ seconds.  Return a lightweight summary instead so
                # generate.py completes in reasonable time.  The entry still
                # gets indexed (with a descriptive name-based summary) and
                # extract_full_text() in agent_base.py handles the same guard
                # so the file won't be opened at query time either.
                MAX_XLSX_BYTES = 50 * 1024 * 1024  # 50 MB
                if file_path.stat().st_size > MAX_XLSX_BYTES:
                    return (
                        f"XLSX (large file — {file_path.stat().st_size // (1024*1024)} MB): "
                        f"{file_path.name}\n"
                        f"This file is too large to index inline. "
                        f"Query it by mentioning the filename or its topic."
                    )

                # ── Smart XLSX snippet ────────────────────────────────────────
                # For large tabular files (revenue, pipeline data) generate a
                # pre-aggregated summary stored in the index entry.
                # extract_full_text() in agent_base.py will serve this cache
                # at query time, avoiding the 50-100s open cost on big files.
                #
                # AGG_KEYWORDS and PREFERRED_NUM_COLS imported from agent_base (module level).
                wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
                text_parts = []

                # ── Sort sheets: largest first so the most data-rich sheet
                # is processed first and dominates the index summary.
                # This fixes files where the active sheet is a tiny pivot
                # (e.g. "1.Whitespace dormant") and the main data is in a
                # sheet called "page" with 400+ rows.
                def _sheet_row_count(s) -> int:
                    try:
                        return s.max_row or 0
                    except Exception:
                        return 0

                sorted_sheets = sorted(
                    wb.worksheets,
                    key=_sheet_row_count,
                    reverse=True,
                )

                for sheet in sorted_sheets:
                    row_iter = sheet.iter_rows(values_only=True)
                    try:
                        hdr = next(row_iter)
                    except StopIteration:
                        continue
                    headers = [str(c) if c is not None else "" for c in hdr]

                    # Probe first 200 rows
                    probe: list = []
                    for row in row_iter:
                        probe.append(row)
                        if len(probe) >= 200:
                            break
                    is_large = len(probe) >= 200

                    text_parts.append(f"[Sheet: {sheet.title}]")

                    if is_large:
                        # Detect numeric column: try preferred names first,
                        # then fall back to last all-numeric col in probe.
                        num_ci: int | None = None
                        h_lower = [h.lower().strip() for h in headers]
                        for pref in PREFERRED_NUM_COLS:
                            if pref in h_lower:
                                ci = h_lower.index(pref)
                                sample = [r[ci] for r in probe
                                          if ci < len(r) and r[ci] is not None]
                                if sample and all(isinstance(v, (int, float)) for v in sample):
                                    num_ci = ci
                                    break
                        if num_ci is None:
                            for ci in range(len(headers) - 1, -1, -1):
                                sample = [r[ci] for r in probe
                                          if ci < len(r) and r[ci] is not None]
                                if sample and all(isinstance(v, (int, float)) for v in sample):
                                    num_ci = ci
                                    break

                        # Detect group-by columns, sorted by AGG_KEYWORDS priority
                        # so that Product appears before Geography in the summary.
                        _agg_order = list(AGG_KEYWORDS.keys())
                        group_cols: list[tuple[int, str]] = sorted(
                            [
                                (ci, AGG_KEYWORDS[h.lower().strip()])
                                for ci, h in enumerate(headers)
                                if h.lower().strip() in AGG_KEYWORDS
                            ],
                            key=lambda x: _agg_order.index(
                                next(k for k in _agg_order if AGG_KEYWORDS[k] == x[1])
                            ),
                        )
                        # Deduplicate: keep only the first occurrence of each label
                        _seen_labels: set[str] = set()
                        group_cols = [
                            (ci, label) for ci, label in group_cols
                            if label not in _seen_labels and not _seen_labels.add(label)  # type: ignore[func-returns-value]
                        ]

                        if num_ci is not None and group_cols:
                            # Single-pass streaming aggregation
                            agg: dict[int, dict[str, float]] = {
                                g_idx: {} for g_idx, _ in group_cols
                            }
                            row_count = len(probe)
                            for row in probe:
                                rev = row[num_ci] if num_ci < len(row) else None
                                if not isinstance(rev, (int, float)):
                                    continue
                                for g_idx, _ in group_cols:
                                    gval = (str(row[g_idx])
                                            if g_idx < len(row) and row[g_idx] is not None
                                            else "(blank)")
                                    agg[g_idx][gval] = agg[g_idx].get(gval, 0.0) + rev
                            for row in row_iter:
                                row_count += 1
                                rev = row[num_ci] if num_ci < len(row) else None
                                if not isinstance(rev, (int, float)):
                                    continue
                                for g_idx, _ in group_cols:
                                    gval = (str(row[g_idx])
                                            if g_idx < len(row) and row[g_idx] is not None
                                            else "(blank)")
                                    agg[g_idx][gval] = agg[g_idx].get(gval, 0.0) + rev

                            text_parts.append(
                                f"Rows: {row_count}  |  Revenue column: '{headers[num_ci]}'\n"
                            )
                            for g_idx, g_label in group_cols:
                                totals = agg[g_idx]
                                if not totals:
                                    continue
                                grand = sum(totals.values())
                                text_parts.append(f"--- By {g_label} ---")
                                for k, v in sorted(totals.items(), key=lambda x: -x[1]):
                                    text_parts.append(f"  {k}: {v:,.2f}")
                                text_parts.append(f"  TOTAL: {grand:,.2f}\n")
                        else:
                            # No aggregatable structure → header + 30-row sample
                            text_parts.append(
                                "Headers: " + " | ".join(h for h in headers if h)
                            )
                            for row in probe[:30]:
                                row_text = " | ".join(str(c) for c in row if c is not None)
                                if row_text.strip():
                                    text_parts.append(row_text)
                    else:
                        # Small sheet → first 30 rows
                        for row in probe[:30]:
                            row_text = " | ".join(str(c) for c in row if c is not None)
                            if row_text.strip():
                                text_parts.append(row_text)

                    text_parts.append("")

                wb.close()
                result = "\n".join(text_parts)
                # If aggregation produced something useful, return it (up to 8000 chars).
                # Otherwise fall back to the plain summary cap.
                cap = max(SUMMARY_CHARS, 8000)
                return result[:cap] if result.strip() else f"XLSX: {file_path.name}"
            except Exception:
                return f"XLSX: {file_path.name}"

        elif ext == ".boxnote":
            try:
                data = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
                def walk(node):
                    if isinstance(node, dict):
                        if node.get("type") == "text":
                            yield node.get("text", "")
                        for v in node.values():
                            yield from walk(v)
                    elif isinstance(node, list):
                        for item in node:
                            yield from walk(item)
                return " ".join(walk(data))[:SUMMARY_CHARS]
            except Exception:
                return f"BoxNote: {file_path.name}"

    except Exception as e:
        return f"{file_path.name} (extraction error: {e})"

    return file_path.name


# ── Index builder ─────────────────────────────────────────────────────────────

def build_index(folder_name: str, force: bool = False) -> dict:
    """
    Build or update the vector index for a single folder.
    Returns { "folder", "entries": [...] }
    Skips unchanged files (hash-based cache).
    """
    folder     = KB_ROOT / folder_name
    idx_path   = index_path_for(folder_name)

    existing: dict[str, dict] = {}
    if idx_path.exists() and not force:
        try:
            data = json.loads(idx_path.read_text())
            existing = {e["path"]: e for e in data.get("entries", [])}
        except Exception:
            existing = {}

    files = [
        f for f in sorted(folder.rglob("*"))
        if f.is_file()
        and f.suffix.lower() in INCLUDE_EXTS
        and not should_skip(f)
    ]

    entries = []
    changed = False

    for f in files:
        rel_path  = str(f.relative_to(KB_ROOT))
        # Use streaming MD5 for large files to avoid loading them all into RAM
        if f.stat().st_size > 10 * 1024 * 1024:  # >10 MB
            h = hashlib.md5()
            with f.open("rb") as fh:
                for chunk in iter(lambda: fh.read(65536), b""):
                    h.update(chunk)
            file_hash = h.hexdigest()
        else:
            file_hash = hashlib.md5(f.read_bytes()).hexdigest()

        if rel_path in existing and existing[rel_path].get("hash") == file_hash:
            entries.append(existing[rel_path])
            continue

        print(f"  Embedding: {f.name}")
        summary    = extract_text_snippet(f)
        embed_text = f"File: {f.name}\n\n{summary}"
        try:
            embedding = get_embedding(embed_text)
        except Exception as e:
            print(f"  Warning: embedding failed for {f.name}: {e}")
            continue

        entries.append({
            "path":      rel_path,
            "name":      f.name,
            "folder":    folder_name,
            "summary":   summary[:2000],
            "embedding": embedding,
            "hash":      file_hash,
        })
        changed = True

    # Remove deleted files
    current_paths = {str(f.relative_to(KB_ROOT)) for f in files}
    entries = [e for e in entries if e["path"] in current_paths]

    if changed or len(entries) != len(existing):
        VECTOR_STORE.mkdir(parents=True, exist_ok=True)
        idx_path.write_text(
            json.dumps({"folder": folder_name, "entries": entries}, indent=2)
        )
        print(f"  ✓ Index saved: {folder_name} ({len(entries)} files)")

    return {"folder": folder_name, "entries": entries}


# ── Semantic search ───────────────────────────────────────────────────────────

def search(query: str, folder_name: str, top_n: int = 4) -> list[dict]:
    """
    Find the top-N most relevant files for a query within a folder.
    Returns list of { path, name, folder, summary, score }.
    Auto-builds the index if missing.
    """
    idx_path = index_path_for(folder_name)

    if not idx_path.exists():
        print(f"  Index not found for {folder_name} — building now...")
        build_index(folder_name)

    data    = json.loads(idx_path.read_text())
    entries = data.get("entries", [])

    if not entries:
        return []

    query_vec = np.array(get_embedding(query)).reshape(1, -1)
    query_dim = query_vec.shape[1]

    # Filter to entries whose embedding dimension matches the query.
    # Mixed dimensions arise when some files were embedded with a different
    # backend (e.g. sentence-transformers=384 vs. Ollama nomic-embed=768).
    # Skipping mismatched entries is safe — they'll be re-indexed on the next
    # generate.py run once a consistent backend is in use.
    compat = [
        e for e in entries
        if isinstance(e.get("embedding"), list)
        and len(e["embedding"]) == query_dim
    ]

    if not compat:
        # No compatible entries — fall back to summary keyword match
        q = query.lower()
        ranked_fb = sorted(
            entries,
            key=lambda e: sum(1 for w in q.split() if w in e.get("summary", "").lower()),
            reverse=True,
        )
        return [
            {
                "path":    e["path"],
                "name":    e["name"],
                "folder":  e["folder"],
                "summary": e.get("summary", ""),
                "score":   0.0,
            }
            for e in ranked_fb[:top_n]
        ]

    doc_vecs  = np.array([e["embedding"] for e in compat])
    scores    = cosine_similarity(query_vec, doc_vecs)[0]

    ranked = sorted(zip(scores, compat), key=lambda x: x[0], reverse=True)
    return [
        {
            "path":    entry["path"],
            "name":    entry["name"],
            "folder":  entry["folder"],
            "summary": entry.get("summary", ""),
            "score":   round(float(score), 4),
        }
        for score, entry in ranked[:top_n]
    ]


def search_all(query: str, top_n: int = 3) -> dict[str, list[dict]]:
    """Search across all discovered folders."""
    return {
        folder: search(query, folder, top_n=top_n)
        for folder in discover_folders()
    }


# ── Per-file index operations (called by watcher) ─────────────────────────────

def update_index_for_file(folder_name: str, file_path: pathlib.Path) -> bool:
    """
    Add or update a single file's embedding in the folder's index.
    Creates the index from scratch if it doesn't exist yet.
    Returns True if the index was modified.
    """
    if file_path.suffix.lower() not in INCLUDE_EXTS or should_skip(file_path):
        return False
    if not file_path.exists():
        return False

    idx_path = index_path_for(folder_name)
    existing_entries: dict[str, dict] = {}
    if idx_path.exists():
        try:
            data = json.loads(idx_path.read_text())
            existing_entries = {e["path"]: e for e in data.get("entries", [])}
        except Exception:
            existing_entries = {}

    rel_path  = str(file_path.relative_to(KB_ROOT))
    if file_path.stat().st_size > 10 * 1024 * 1024:
        h = hashlib.md5()
        with file_path.open("rb") as fh:
            for chunk in iter(lambda: fh.read(65536), b""):
                h.update(chunk)
        fhash = h.hexdigest()
    else:
        fhash = hashlib.md5(file_path.read_bytes()).hexdigest()

    # Skip if already indexed at the same hash
    if rel_path in existing_entries and existing_entries[rel_path].get("hash") == fhash:
        return False

    snippet    = extract_text_snippet(file_path)
    embed_text = f"File: {file_path.name}\n\n{snippet}"
    try:
        embedding = get_embedding(embed_text)
    except Exception as e:
        print(f"  Warning: embedding failed for {file_path.name}: {e}")
        return False

    existing_entries[rel_path] = {
        "path":      rel_path,
        "name":      file_path.name,
        "folder":    folder_name,
        "summary":   snippet[:2000],
        "embedding": embedding,
        "hash":      fhash,
    }

    VECTOR_STORE.mkdir(parents=True, exist_ok=True)
    idx_path.write_text(
        json.dumps({"folder": folder_name, "entries": list(existing_entries.values())}, indent=2)
    )
    return True


def remove_from_index(folder_name: str, file_path: pathlib.Path) -> bool:
    """
    Remove a single file's entry from the folder's index.
    Returns True if the index was modified.
    """
    idx_path = index_path_for(folder_name)
    if not idx_path.exists():
        return False

    try:
        data    = json.loads(idx_path.read_text())
        entries = data.get("entries", [])
    except Exception:
        return False

    rel_path = str(file_path.relative_to(KB_ROOT))
    new_entries = [e for e in entries if e["path"] != rel_path]

    if len(new_entries) == len(entries):
        return False  # nothing removed

    idx_path.write_text(
        json.dumps({"folder": folder_name, "entries": new_entries}, indent=2)
    )
    return True


def delete_index(folder_name: str):
    """Delete the entire index file for a folder (called when folder is removed)."""
    idx_path = index_path_for(folder_name)
    if idx_path.exists():
        idx_path.unlink()


# ── Standalone: rebuild all indexes ──────────────────────────────────────────

if __name__ == "__main__":
    import sys
    force = "--force" in sys.argv

    folders = discover_folders()
    if not folders:
        print(f"No knowledge folders found under {KB_ROOT}")
        sys.exit(1)

    print(f"KnowledgeBase — Building vector indexes (KB_ROOT={KB_ROOT})\n")
    for folder_name in folders:
        print(f"Indexing {folder_name}...")
        index = build_index(folder_name, force=force)
        print(f"  → {len(index['entries'])} files indexed\n")

    print("Done ✓")
