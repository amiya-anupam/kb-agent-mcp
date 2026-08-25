#!/usr/bin/env python3
"""
agent_base.py — Shared base logic for all KnowledgeBase sub-agents
-------------------------------------------------------------------
Generated agents (agent_<folder>.py) are thin wrappers that set their
folder-specific config and delegate everything here.

Provides:
  - extract_full_text(file_path)  — multi-format text extractor
  - call_llm(messages)            — provider-agnostic LLM call
  - ask(question, folder, ...)    — README-first RAG pipeline

Supported formats in extract_full_text()
-----------------------------------------
.txt / .md / .csv       plain read
.docx                   XML extraction (zipfile)
.pdf                    pypdf page iteration
.pptx / .ppt            python-pptx shape text
.xlsx / .xls            streaming aggregation or openpyxl
.boxnote                JSON tree walk
.png / .jpg / .jpeg     OCR via pytesseract (KB_OCR_ENABLED=true, default)
.gif / .webp              or PIL metadata fallback, or filename-only

README-first pipeline:
  1. Find the folder's README (must contain <!-- KB:AUTO-INDEX:START --> block)
  2. For normal questions  → pass AUTO-INDEX block only (~500-2000 tokens)
  3. For complex questions → pass full README up to KB_BUDGET_FULL_README chars
  4. Fallback              → raw-file RAG if README is absent or too thin

LLM provider is driven entirely by env vars:
  KB_LLM_PROVIDER  ollama | openai | anthropic | custom | passthrough
  KB_LLM_BASE_URL  base endpoint
  KB_MODEL         model name
  KB_API_KEY       API key (openai / anthropic / custom)

Passthrough mode (KB_LLM_PROVIDER=passthrough, or auto-detected when no
LLM is reachable and KB_PASSTHROUGH_FALLBACK != "false"):
  Instead of calling a local LLM, the agent emits a structured
  <<<KB_PASSTHROUGH>>> block to stdout. Bob's Claude reads that block and
  answers directly — no local LLM required.
"""

import os
import re
import json
import pathlib
import sys

# ── Environment loader ────────────────────────────────────────────────────────

def _load_env():
    for candidate in [
        pathlib.Path(os.environ.get("KB_ROOT", "")) / ".env",
        pathlib.Path(__file__).parent.parent / ".env",
    ]:
        if candidate.exists():
            for line in candidate.read_text(encoding="utf-8").splitlines():
                line = line.strip()
                if line and not line.startswith("#") and "=" in line:
                    k, _, v = line.partition("=")
                    os.environ.setdefault(k.strip(), v.strip())
            break

_load_env()

# ── Config (all from env) ─────────────────────────────────────────────────────

def _kb_root() -> pathlib.Path:
    raw = os.environ.get("KB_ROOT", "")
    if raw:
        return pathlib.Path(raw)
    return pathlib.Path(__file__).parent.parent

KB_ROOT      = _kb_root()
VECTOR_STORE = pathlib.Path(__file__).parent / "vector_store"

# ── XLSX aggregation config ───────────────────────────────────────────────────
# Single source of truth for both _stream_xlsx_aggregate() (>50 MB path) and
# extract_full_text() (<50 MB path).  embeddings.py imports these directly.
# kb_agent_mcp/file_parser.py keeps its own copy (_AGG_KEYWORDS) — intentional
# package-boundary isolation; do not remove that copy.

AGG_KEYWORDS: dict[str, str] = {
    # ── Product columns first (most semantically relevant for search) ──
    "ut lvl 30 name dynamic": "Product (UT L30)",
    "ut l30 name":            "Product (UT L30)",
    "ut l30":                 "Product (UT L30)",
    "product family name":    "Product Family",
    "reporting product family": "Reporting Product Family",
    "product":                "Product",
    # ── Time columns ──
    "year": "Year", "quarter": "Quarter",
    "quarter in year": "Quarter In Year",
    # ── Standard finance columns ──
    "geography":                    "Geography",
    "geography name":               "Geography",
    "market":                       "Market",
    "market name":                  "Market",
    "country":                      "Country",
    "finance family":               "Finance Family",
    "revenue type":                 "Revenue Type",
    "reporting revenue type name":  "Revenue Type",
    "on-prem or saas":              "On-prem/SaaS",
    "division":                     "Division",
    # ── CRM deal columns ──
    "classification name":          "Classification",
    "frozen client lifecycle name": "Client Lifecycle",
    # ── Renewal / ELA columns ──
    "status": "Status",
}

# Preferred numeric column names — checked in order; first match wins.
# CRM exports use "Won"; revenue reports use "Rev Act @ PC".
PREFERRED_NUM_COLS: list[str] = [
    "won", "total(cy cw won @ pc)", "rev act @ pc",
    "amount", "oppty value", "total",
]


# ── File discovery constants ──────────────────────────────────────────────────
# Single source of truth for domain-discovery in embeddings.py, generate.py,
# and watch_kb.py.  watch_kb.py intentionally *extends* INCLUDE_EXTS with
# image types and uses a different SKIP_PATTERNS (must exclude .watch.log,
# thumbs.db, ~$ temp files) — it imports DEFAULT_BLOCKLIST only.
# ask.py and agent_knowledgebase.py have larger blocklists (they add infra
# folders like scripts/, tests/, .kb_index/) and keep their own local copies.

DEFAULT_BLOCKLIST: frozenset[str] = frozenset({
    "agents", ".git", "__pycache__", ".ds_store", "node_modules",
    ".venv", "venv", "env", ".bob", ".idea", ".vscode", "dist", "build",
})

INCLUDE_EXTS: frozenset[str] = frozenset({
    ".pdf", ".docx", ".pptx", ".xlsx", ".md", ".txt",
    ".csv", ".boxnote", ".ppt", ".doc",
})

SKIP_PATTERNS: frozenset[str] = frozenset({
    "readme", ".ds_store", "watch_kb", "__pycache__",
})


# ── .noindex sentinel check ───────────────────────────────────────────────────

def _has_noindex_ancestor(path: pathlib.Path) -> bool:
    """
    Return True if any ancestor directory of *path* (up to KB_ROOT) contains
    a `.noindex` sentinel file.

    A `.noindex` file placed in a folder (or any of its parents) permanently
    excludes all files beneath it from being read, indexed, or returned in
    RAG context.  This is the agents-layer equivalent of the canonical
    implementation in kb_agent_mcp/file_parser._has_noindex_ancestor().

    The check walks upward from the file's immediate parent and stops once it
    reaches KB_ROOT itself, so it never escapes the knowledge-base boundary.
    It is best-effort — any I/O error is silently swallowed so the guard
    never crashes a query in progress.
    """
    try:
        for parent in path.parents:
            if (parent / ".noindex").exists():
                return True
            # Stop after checking KB_ROOT itself — don't walk the whole filesystem
            if parent == KB_ROOT:
                break
    except Exception:
        pass
    return False

def folder_to_safe_name(name: str) -> str:
    """Convert a folder name to a safe snake_case identifier (e.g. 'ACE Docs' → 'ace_docs')."""
    n = name.lower()
    n = re.sub(r"[^a-z0-9]+", "_", n)
    return n.strip("_")


def should_skip(path: pathlib.Path) -> bool:
    """Return True if this file should be excluded from indexing.

    Skips:
    • files matching SKIP_PATTERNS (readme, .ds_store, etc.)
    • any file whose ancestor folder contains a `.noindex` sentinel file
    """
    name = path.name.lower()
    if any(p in name for p in SKIP_PATTERNS):
        return True
    return _has_noindex_ancestor(path)


def _apply_format_instruction(system_prompt: str, format_instruction: str) -> str:
    """
    Append a format instruction to a system prompt.

    Defined here (agent_base) so per-domain sub-agent files can import it
    without importing from agent_knowledgebase (which would create a circular
    dependency).  agent_knowledgebase imports and re-exports this function.
    """
    if not format_instruction:
        return system_prompt
    return system_prompt + f"\n\n**OUTPUT FORMAT DIRECTIVE (highest priority):**\n{format_instruction}"

LLM_PROVIDER = os.environ.get("KB_LLM_PROVIDER", "ollama").lower()
LLM_BASE_URL = os.environ.get("KB_LLM_BASE_URL", "http://localhost:11434")
MODEL        = os.environ.get("KB_MODEL", "qwen3:14b")
API_KEY      = os.environ.get("KB_API_KEY", "")

# ── Passthrough helpers ───────────────────────────────────────────────────────

_PASSTHROUGH_MARKER = "<<<KB_PASSTHROUGH>>>"
_PASSTHROUGH_END    = "<<<KB_PASSTHROUGH_END>>>"

def _is_passthrough_mode() -> bool:
    """
    Return True when the agent should emit a passthrough block instead of
    calling a local LLM.

    Conditions (any of):
      1. KB_LLM_PROVIDER is explicitly set to "passthrough"
      2. KB_LLM_PROVIDER is "ollama" (default), no KB_API_KEY is set,
         and the Ollama endpoint is not reachable
         (auto-detection; can be disabled with KB_PASSTHROUGH_FALLBACK=false)
    """
    if LLM_PROVIDER == "passthrough":
        return True

    # Only auto-detect for the default Ollama provider when no API key is set
    if LLM_PROVIDER not in ("ollama",) or API_KEY:
        return False

    if os.environ.get("KB_PASSTHROUGH_FALLBACK", "true").lower() == "false":
        return False

    try:
        import httpx
        r = httpx.get(f"{LLM_BASE_URL}/api/tags", timeout=3.0)
        return r.status_code >= 400
    except Exception:
        return True  # unreachable → passthrough


# Evaluated once at import so sub-agents share the same decision
_PASSTHROUGH = _is_passthrough_mode()


def emit_passthrough(question: str, context: str, system_prompt: str,
                     agent_name: str, source_label: str) -> dict:
    """
    Print a structured passthrough block to stdout and return a sentinel dict.

    Bob's skill handler reads everything the script prints.  The block is
    human-readable so Bob's Claude can parse it without any special tooling:

        <<<KB_PASSTHROUGH>>>
        AGENT: <name>
        QUESTION: <question>
        SOURCE: <label>
        SYSTEM_PROMPT:
        <system prompt>
        ---CONTEXT---
        <retrieved context>
        <<<KB_PASSTHROUGH_END>>>

    Bob's Claude sees this output and answers the question using the context
    provided, then returns the answer to the user.
    """
    block = (
        f"\n{_PASSTHROUGH_MARKER}\n"
        f"AGENT: {agent_name}\n"
        f"QUESTION: {question}\n"
        f"SOURCE: {source_label}\n"
        f"SYSTEM_PROMPT:\n{system_prompt}\n"
        f"---CONTEXT---\n{context}\n"
        f"{_PASSTHROUGH_END}\n"
    )
    print(block, flush=True)
    return {
        "agent":   agent_name,
        "answer":  block,   # orchestrator treats this as the answer text
        "sources": [{"name": source_label, "path": source_label, "score": 1.0}],
        "found":   True,
        "passthrough": True,
    }

# README-first config
MARKER_START     = "<!-- KB:AUTO-INDEX:START -->"
MARKER_END       = "<!-- KB:AUTO-INDEX:END -->"

import importlib as _importlib
_cb = _importlib.import_module("context_budget")

# Keywords that indicate the user genuinely needs the full README narrative.
#
# INTENTIONALLY NARROW: casual phrasing ("tell me about", "describe",
# "overview of") is NOT included — those questions are almost always
# answered by the compacted AUTO-INDEX block (index mode, ~2 000 tokens).
# Only structural/comparative questions that need the full document body
# are included here.  Incorrect inclusion of casual verbs was the single
# largest source of token waste (up to 6 000 tokens per domain per query).
_COMPLEX_QUESTION_PATTERNS = re.compile(
    r"\b(compare|contrast|difference between|differences between|"
    r"walk me through|step[- ]by[- ]step|deep dive|in[- ]depth|"
    r"comprehensive|explain in detail|elaborate on|how does .{3,40} work|"
    r"pros and cons|trade[- ]off|architecture of|internals of|"
    r"full breakdown|everything about)\b",
    re.IGNORECASE,
)

# Questions that require reading actual file data — the README AUTO-INDEX only
# stores short summaries, so these must bypass README-first and use raw-file RAG
# which calls extract_full_text() on the matched files directly.
_DATA_QUESTION_PATTERNS = re.compile(
    r"\b(revenue|total revenue|arr|mrr|acv|tcv|quota|attainment|"
    r"how much|how many|what is the (number|count|total|sum|amount|value|"
    r"figure|balance|price|cost|rate|percentage|percent|ratio|score|metric)|"
    r"what (are|were) the (number|count|total|sum|figures|numbers|metrics|results|"
    r"revenue|sales|deals|renewals|bookings|customers|accounts)|"
    r"list (all|every|the|each)|show me (all|the|every)|give me (all|the)|"
    r"breakdown|by (quarter|region|country|geography|market|segment|"
    r"product|customer|account|industry|channel)|"
    r"q[1-4]\s*20\d\d|20\d\d\s*q[1-4]|fy\s*20\d\d|"
    r"ytd|yoy|qoq|mom|r4q|trailing|rolling)\b",
    re.IGNORECASE,
)


def _is_data_question(question: str) -> bool:
    """Return True when the question asks for concrete data, numbers, or lists."""
    return bool(_DATA_QUESTION_PATTERNS.search(question))


# ── README helpers ────────────────────────────────────────────────────────────

def _find_readme(folder: pathlib.Path) -> pathlib.Path | None:
    """
    Locate the README for a knowledge folder using a priority cascade:
      1. Any .md whose name contains 'readme' (case-insensitive)
      2. <FolderName>.md  (standard name used by generate.py)
      3. Any .md file whose first 500 chars contain a Markdown heading (# …)
      4. The first .md file found (last resort)

    This makes README discovery fully dynamic — it works regardless of what
    the user named the file (e.g. 'ACE readme file.md', 'ACE Docs.md',
    'Jon Doe Analytics Overview.md', etc.).
    """
    try:
        md_files = [f for f in folder.iterdir()
                    if f.is_file() and f.suffix.lower() == ".md"]
    except Exception:
        return None

    if not md_files:
        return None

    # Priority 1: name contains "readme"
    for f in md_files:
        if "readme" in f.name.lower():
            return f

    # Priority 2: matches the folder name exactly (e.g. "ACE Docs.md")
    folder_name_md = folder.name + ".md"
    for f in md_files:
        if f.name == folder_name_md:
            return f

    # Priority 3: first .md whose content starts with a Markdown heading
    for f in md_files:
        try:
            head = f.read_text(encoding="utf-8", errors="ignore")[:500]
            if re.search(r"^#{1,3}\s+\S", head, re.MULTILINE):
                return f
        except Exception:
            continue

    # Priority 4: first .md file
    return md_files[0]


def _extract_auto_index_block(readme_text: str) -> str | None:
    """Extract text between AUTO-INDEX markers. Returns None if markers absent."""
    if MARKER_START not in readme_text or MARKER_END not in readme_text:
        return None
    start = readme_text.index(MARKER_START) + len(MARKER_START)
    end   = readme_text.index(MARKER_END)
    return readme_text[start:end].strip()


def _non_index_chars(readme_text: str) -> int:
    """Count characters in the README outside the AUTO-INDEX block."""
    if MARKER_START in readme_text and MARKER_END in readme_text:
        start = readme_text.index(MARKER_START)
        end   = readme_text.index(MARKER_END) + len(MARKER_END)
        outside = readme_text[:start] + readme_text[end:]
    else:
        outside = readme_text
    return len(outside.strip())


def _is_complex_question(question: str) -> bool:
    return bool(_COMPLEX_QUESTION_PATTERNS.search(question))


def _get_readme_context(folder_name: str, question: str) -> tuple[str | None, str]:
    """
    Return (context_text, source_label) for the README-first strategy.

    Returns (None, "") if README is absent, too thin, or the question is a
    data/numeric question — caller should fall back to raw-file RAG.
    """
    # Data questions need real file content, not README summaries.
    # Short-circuit immediately so raw-file RAG opens and reads the actual files.
    if _is_data_question(question):
        return None, ""

    folder = KB_ROOT / folder_name
    readme = _find_readme(folder)

    if readme is None:
        return None, ""

    try:
        readme_text = readme.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return None, ""

    # README must have meaningful hand-written content (not just the auto-index)
    if _non_index_chars(readme_text) < _cb.get("min_readme"):
        return None, ""

    is_complex = _is_complex_question(question)

    if is_complex:
        # Complex question: full README body, compacted and capped
        context = _cb.trim(readme_text, "full_readme")
        label   = f"Full README ({readme.name})"
    else:
        # Simple question: compact AUTO-INDEX block + brief intro
        auto_index = _extract_auto_index_block(readme_text)
        if auto_index:
            pre_index = readme_text[:readme_text.index(MARKER_START)].strip()
            context   = _cb.build_context(pre_index, auto_index)
            label     = f"README index ({readme.name})"
        else:
            # README has no auto-index block yet — use full README
            context = _cb.trim(readme_text, "full_readme")
            label   = f"Full README ({readme.name})"

    return context, label



# ── Streaming aggregator for very large XLSX files ────────────────────────────
# _stream_xlsx_aggregate is the canonical XLSX streaming aggregation entry point.
# The pure algorithm lives in _xlsx_stream.py; this wrapper binds the agents-layer
# AGG_KEYWORDS and PREFERRED_NUM_COLS constants.

import _xlsx_stream as _xlsx_stream_mod  # noqa: E402

def _stream_xlsx_aggregate(file_path: pathlib.Path, max_chars: int = 8000) -> str:
    """Stream-aggregate a large XLSX file. Delegates to the shared _xlsx_stream module."""
    return _xlsx_stream_mod.stream_xlsx_aggregate(
        file_path, max_chars, AGG_KEYWORDS, PREFERRED_NUM_COLS
    )


# ── Full text extractor ───────────────────────────────────────────────────────

def extract_full_text(file_path: pathlib.Path, max_chars: int | None = None) -> str:
    """Extract as much useful text as possible from a file."""
    if max_chars is None:
        max_chars = _cb.get("rag_file")

    # ── Security: honour .noindex sentinels ──────────────────────────────────
    # Any file whose ancestor folder contains a `.noindex` file is hard-excluded.
    # This mirrors the enforcement in kb_agent_mcp/file_parser.should_skip() and
    # kb_agent_mcp/security_gate so the agents-layer cannot be used as a bypass
    # path to read protected content.
    if _has_noindex_ancestor(file_path):
        return f"[Excluded: {file_path.name} is in a .noindex protected folder]"

    ext = file_path.suffix.lower()
    try:
        if ext in {".txt", ".md", ".csv"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")[:max_chars]

        elif ext == ".docx":
            import zipfile, re as _re
            with zipfile.ZipFile(file_path) as z:
                with z.open("word/document.xml") as f:
                    xml = f.read().decode("utf-8", errors="ignore")
            text = _re.sub(r"<[^>]+>", " ", xml)
            return " ".join(text.split())[:max_chars]

        elif ext == ".pdf":
            try:
                from pypdf import PdfReader
                reader = PdfReader(str(file_path))
                text = ""
                for page in reader.pages:
                    text += (page.extract_text() or "") + "\n"
                    if len(text) >= max_chars:
                        break
                return text[:max_chars]
            except Exception as e:
                return f"[PDF read error: {e}]"

        elif ext in {".pptx", ".ppt"}:
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text.strip() + "\n"
                    if len(text) >= max_chars:
                        break
                return text[:max_chars]
            except Exception as e:
                return f"[PPTX read error: {e}]"

        elif ext in {".xlsx", ".xls"}:
            try:
                import openpyxl

                # ── Cache-first XLSX extraction ───────────────────────────────
                # Serve from the pre-aggregated summary stored in the vector
                # index by build_index().  If the cache has a rich summary
                # (>200 chars), return it directly.  This MUST run before the
                # >50 MB guard so large files skip the 50s streaming parse when
                # a cached summary is already available.
                idx_path = VECTOR_STORE / f"{folder_to_safe_name(file_path.parent.name)}_index.json"
                if idx_path.exists():
                    try:
                        idx_data = json.loads(idx_path.read_text())
                        rel_path = str(file_path.relative_to(KB_ROOT))
                        for entry in idx_data.get("entries", []):
                            if entry.get("path") == rel_path:
                                cached = entry.get("summary", "")
                                if len(cached) > 200:
                                    return cached[:max_chars]
                                break
                    except Exception:
                        pass  # cache miss → fall through to live read

                # ── Guard: very large XLSX files (>50 MB) ─────────────────────
                # openpyxl takes 80+ seconds to open a 160 MB file. For such
                # files we use a fast streaming XML+iterparse approach instead.
                MAX_XLSX_BYTES = 50 * 1024 * 1024  # 50 MB
                if file_path.stat().st_size > MAX_XLSX_BYTES:
                    return _stream_xlsx_aggregate(file_path, max_chars)

                # ── Live single-pass streaming XLSX extraction ────────────────
                # For large tabular files (revenue, pipeline data) we aggregate
                # into per-dimension totals in a single streaming pass instead of
                # materialising all rows.  The file is opened exactly once.
                #
                # Algorithm:
                #   1. Read header row → detect numeric column (preferred by
                #      well-known names, else last all-numeric col in probe)
                #      and categorical group-by cols by matching header names.
                #   2. Continue streaming the rest of the file, accumulating
                #      per-group totals on the fly.
                #   3. Emit aggregated summary.  Falls back to first-200-rows dump
                #      when no aggregatable structure is found.

                text_parts = []
                wb = openpyxl.load_workbook(
                    str(file_path), read_only=True, data_only=True
                )

                # ── Sort sheets: largest first so the most data-rich sheet
                # is processed first and dominates the context output.
                def _sheet_row_count(s) -> int:
                    try:
                        return s.max_row or 0
                    except Exception:
                        return 0

                sorted_sheets = sorted(
                    wb.worksheets,
                    key=_sheet_row_count,
                    reverse=True,
                )

                for sheet in sorted_sheets:
                    row_iter = sheet.iter_rows(values_only=True)

                    # ── Header row ────────────────────────────────────────────
                    try:
                        hdr = next(row_iter)
                    except StopIteration:
                        continue
                    headers = [str(c) if c is not None else "" for c in hdr]

                    # ── Probe: first 200 rows to detect column types ───────────
                    probe: list = []
                    for row in row_iter:
                        probe.append(row)
                        if len(probe) >= 200:
                            break
                    is_large = len(probe) >= 200

                    text_parts.append(f"[Sheet: {sheet.title}]")

                    if is_large:
                        # Detect numeric column: try preferred names first,
                        # then fall back to last all-numeric col in probe.
                        num_ci: int | None = None
                        h_lower = [h.lower().strip() for h in headers]
                        for pref in PREFERRED_NUM_COLS:
                            if pref in h_lower:
                                ci = h_lower.index(pref)
                                sample = [r[ci] for r in probe
                                          if ci < len(r) and r[ci] is not None]
                                if sample and all(isinstance(v, (int, float)) for v in sample):
                                    num_ci = ci
                                    break
                        if num_ci is None:
                            for ci in range(len(headers) - 1, -1, -1):
                                sample = [r[ci] for r in probe
                                          if ci < len(r) and r[ci] is not None]
                                if sample and all(isinstance(v, (int, float)) for v in sample):
                                    num_ci = ci
                                    break

                        # Detect group-by columns, sorted by AGG_KEYWORDS priority
                        _agg_order = list(AGG_KEYWORDS.keys())
                        group_cols: list[tuple[int, str]] = sorted(
                            [
                                (ci, AGG_KEYWORDS[h.lower().strip()])
                                for ci, h in enumerate(headers)
                                if h.lower().strip() in AGG_KEYWORDS
                            ],
                            key=lambda x: _agg_order.index(
                                next(k for k in _agg_order if AGG_KEYWORDS[k] == x[1])
                            ),
                        )
                        _seen_labels: set[str] = set()
                        group_cols = [
                            (ci, label) for ci, label in group_cols
                            if label not in _seen_labels and not _seen_labels.add(label)  # type: ignore[func-returns-value]
                        ]

                        if num_ci is not None and group_cols:
                            # ── Streaming aggregation (single pass) ───────────
                            agg: dict[int, dict[str, float]] = {
                                g_idx: {} for g_idx, _ in group_cols
                            }
                            row_count = len(probe)

                            # Accumulate probe rows first
                            for row in probe:
                                rev = row[num_ci] if num_ci < len(row) else None
                                if not isinstance(rev, (int, float)):
                                    continue
                                for g_idx, _ in group_cols:
                                    gval = (str(row[g_idx])
                                            if g_idx < len(row) and row[g_idx] is not None
                                            else "(blank)")
                                    agg[g_idx][gval] = agg[g_idx].get(gval, 0.0) + rev

                            # Continue streaming remaining rows
                            for row in row_iter:
                                row_count += 1
                                rev = row[num_ci] if num_ci < len(row) else None
                                if not isinstance(rev, (int, float)):
                                    continue
                                for g_idx, _ in group_cols:
                                    gval = (str(row[g_idx])
                                            if g_idx < len(row) and row[g_idx] is not None
                                            else "(blank)")
                                    agg[g_idx][gval] = agg[g_idx].get(gval, 0.0) + rev

                            text_parts.append(
                                f"Rows: {row_count}  |  Revenue column: '{headers[num_ci]}'\n"
                            )
                            for g_idx, g_label in group_cols:
                                totals = agg[g_idx]
                                if not totals:
                                    continue
                                grand = sum(totals.values())
                                text_parts.append(f"--- By {g_label} ---")
                                for k, v in sorted(totals.items(), key=lambda x: -x[1]):
                                    text_parts.append(f"  {k}: {v:,.2f}")
                                text_parts.append(f"  TOTAL: {grand:,.2f}\n")
                        else:
                            # Large but no aggregatable structure → probe sample
                            text_parts.append(
                                f"Rows: 200+ (sample)  Columns: {len(headers)}"
                            )
                            text_parts.append(
                                "Headers: " + " | ".join(h for h in headers if h) + "\n"
                            )
                            for row in probe[:100]:
                                row_text = " | ".join(
                                    str(c) for c in row if c is not None
                                )
                                if row_text.strip():
                                    text_parts.append(row_text)
                    else:
                        # Small sheet — dump all rows verbatim
                        text_parts.append(
                            f"Rows: {len(probe)}  Columns: {len(headers)}"
                        )
                        text_parts.append(
                            "Headers: " + " | ".join(h for h in headers if h) + "\n"
                        )
                        for row in probe:
                            row_text = " | ".join(
                                str(c) for c in row if c is not None
                            )
                            if row_text.strip():
                                text_parts.append(row_text)

                    text_parts.append("")  # blank line between sheets

                wb.close()
                return "\n".join(text_parts)[:max_chars]
            except Exception as e:
                return f"[XLSX read error: {e}]"

        elif ext == ".boxnote":
            try:
                data = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
                def walk(node):
                    if isinstance(node, dict):
                        if node.get("type") == "text":
                            yield node.get("text", "")
                        for v in node.values():
                            yield from walk(v)
                    elif isinstance(node, list):
                        for item in node:
                            yield from walk(item)
                return " ".join(walk(data))[:max_chars]
            except Exception as e:
                return f"[BoxNote read error: {e}]"

        elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
            return _extract_image_text(file_path, max_chars)

    except Exception as e:
        return f"[Read error: {e}]"

    return f"[Unsupported: {file_path.name}]"


# ── Image OCR helper (agents-layer) ──────────────────────────────────────────

def _extract_image_text(path: pathlib.Path, max_chars: int) -> str:
    """
    Extract text from an image file.

    Mirrors the logic in kb_agent_mcp/file_parser._extract_image().

    Strategy (in priority order):
    1. pytesseract OCR  — when KB_OCR_ENABLED != "false" (default: enabled)
                          and KB_OCR_ENGINE in ("tesseract", "auto") (default).
    2. PIL/Pillow metadata fallback — image dimensions + filename.
    3. Filename-only last resort.

    All optional deps (pytesseract, PIL) are imported lazily so the agents
    layer works without them installed — it just falls back to lower tiers.
    """
    ocr_enabled = os.environ.get("KB_OCR_ENABLED", "true").lower() not in ("false", "0", "no")
    if not ocr_enabled:
        return f"[Image: {path.name}]"

    ocr_engine = os.environ.get("KB_OCR_ENGINE", "auto").lower()

    # ── Try pytesseract ────────────────────────────────────────────────────────
    if ocr_engine in ("tesseract", "auto"):
        try:
            import pytesseract
            from PIL import Image as _Image
            img = _Image.open(str(path))
            text = pytesseract.image_to_string(img).strip()
            if text:
                return text[:max_chars]
        except ImportError:
            pass  # pytesseract / PIL not installed — fall through to PIL-only
        except Exception as exc:
            import logging as _logging
            _logging.getLogger(__name__).debug(
                "pytesseract failed for %s: %s", path.name, exc
            )

    # ── PIL metadata fallback ──────────────────────────────────────────────────
    try:
        from PIL import Image as _Image
        img = _Image.open(str(path))
        w, h = img.size
        mode = img.mode
        return f"[Image: {path.name} | {w}×{h}px | mode={mode}]"[:max_chars]
    except ImportError:
        pass
    except Exception as exc:
        import logging as _logging
        _logging.getLogger(__name__).debug("PIL failed for %s: %s", path.name, exc)

    return f"[Image: {path.name}]"


# ── LLM call (provider-agnostic) ──────────────────────────────────────────────

def call_llm(messages: list[dict], temperature: float = 0.2) -> str:
    """
    Send messages to the configured LLM and return the response text.
    Supports: ollama | openai | anthropic | custom
    """
    if LLM_PROVIDER == "anthropic":
        return _call_anthropic(messages, temperature)

    if LLM_PROVIDER in ("openai", "custom"):
        return _call_openai_compat(messages, temperature)

    # Default: Ollama
    return _call_ollama(messages, temperature)


def call_llm_generate(prompt: str) -> str:
    """
    Call the LLM with a plain text prompt and return the response.
    Convenience wrapper used by generate.py and watch_kb.py for single-turn
    generation tasks (domain descriptions, file summaries, etc.).
    Uses temperature=0.3 — slightly more creative than the default RAG setting.
    """
    return call_llm([{"role": "user", "content": prompt}], temperature=0.3)


def _call_ollama(messages: list[dict], temperature: float) -> str:
    import httpx
    try:
        response = httpx.post(
            f"{LLM_BASE_URL}/api/chat",
            json={
                "model":    MODEL,
                "messages": messages,
                "stream":   False,
                "options":  {"temperature": temperature, "num_ctx": _cb.get("num_ctx")},
                "think":    False,
            },
            timeout=120.0,
        )
        response.raise_for_status()
        return response.json()["message"]["content"].strip()
    except Exception as e:
        hint = (
            f"Ollama call failed: {e}\n"
            f"  URL:   {LLM_BASE_URL}/api/chat\n"
            f"  Model: {MODEL}\n"
            f"  Check: is Ollama running? (`ollama serve`)\n"
            f"         is the model pulled? (`ollama pull {MODEL}`)"
        )
        raise RuntimeError(hint) from e


def _call_openai_compat(messages: list[dict], temperature: float) -> str:
    import httpx
    base = LLM_BASE_URL.rstrip("/")
    if "11434" in base and not base.endswith("/v1"):
        base = f"{base}/v1"
    headers = {"Content-Type": "application/json"}
    if API_KEY:
        headers["Authorization"] = f"Bearer {API_KEY}"
    try:
        response = httpx.post(
            f"{base}/chat/completions",
            headers=headers,
            json={
                "model":       MODEL,
                "messages":    messages,
                "temperature": temperature,
            },
            timeout=120.0,
        )
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"].strip()
    except Exception as e:
        missing_key = not API_KEY and LLM_PROVIDER in ("openai", "custom")
        hint = (
            f"OpenAI-compatible call failed: {e}\n"
            f"  URL:   {base}/chat/completions\n"
            f"  Model: {MODEL}\n"
        )
        if missing_key:
            hint += "  Check: KB_API_KEY is not set in your .env\n"
        else:
            hint += "  Check: is KB_API_KEY correct? Is the endpoint reachable?\n"
        raise RuntimeError(hint) from e


def _call_anthropic(messages: list[dict], temperature: float) -> str:
    import httpx
    system      = ""
    chat_messages = []
    for m in messages:
        if m["role"] == "system":
            system = m["content"]
        else:
            chat_messages.append(m)

    headers = {
        "x-api-key":         API_KEY,
        "anthropic-version": "2023-06-01",
        "Content-Type":      "application/json",
    }
    payload = {
        "model":       MODEL,
        "max_tokens":  4096,
        "temperature": temperature,
        "messages":    chat_messages,
    }
    if system:
        payload["system"] = system

    try:
        response = httpx.post(
            f"{LLM_BASE_URL}/v1/messages",
            headers=headers,
            json=payload,
            timeout=120.0,
        )
        response.raise_for_status()
        return response.json()["content"][0]["text"].strip()
    except Exception as e:
        hint = (
            f"Anthropic call failed: {e}\n"
            f"  URL:   {LLM_BASE_URL}/v1/messages\n"
            f"  Model: {MODEL}\n"
            f"  Check: is KB_API_KEY set correctly in your .env?\n"
            f"         does your key have access to {MODEL}?"
        )
        raise RuntimeError(hint) from e


# ── Confidence footer ─────────────────────────────────────────────────────────

def format_confidence_footer(sources: list[dict]) -> str:
    """
    Build a human-readable confidence footer from a sources list.

    Uses the top source's score (0.0–1.0 cosine similarity) to assign a
    label and returns a single footer line, e.g.:

        Confidence: High (0.87) — Source: Q3_Renewal_Tracker.xlsx

    For README-first results the score is 1.0 and the label is omitted
    (the source IS the authoritative document — no retrieval uncertainty).

    Score thresholds:
        ≥ 0.80  → High
        ≥ 0.60  → Medium
        < 0.60  → Low
    """
    if not sources:
        return ""

    top = sources[0]
    score = top.get("score", 0.0)
    name  = top.get("name", "unknown")

    # README-first path sets score=1.0 to signal "authoritative, no vector search"
    # Show source without a confidence label in that case.
    if score >= 1.0:
        return f"\n\n---\n📄 **Source:** `{name}`"

    if score >= 0.80:
        label = "High"
    elif score >= 0.60:
        label = "Medium"
    else:
        label = "Low"

    # Include up to 2 additional sources if present
    extra = sources[1:3]
    extra_str = ""
    if extra:
        extra_str = " · " + " · ".join(f"`{s['name']}`" for s in extra)

    return (
        f"\n\n---\n🎯 **Confidence:** {label} ({score:.2f})"
        f" — **Source:** `{name}`{extra_str}"
    )


# ── Core ask function ─────────────────────────────────────────────────────────

def ask(
    question: str,
    folder_name: str,
    agent_name: str,
    system_prompt: str,
    conversation_history: list[dict] | None = None,
    top_n: int = 4,
    max_chars: int | None = None,
    _pre_ranked_results: list[dict] | None = None,
) -> dict:
    """
    README-first RAG pipeline for a single folder domain.

    Strategy:
      1. Try README-first: use the folder README as primary context
         - Normal questions  → AUTO-INDEX block + brief intro section
         - Complex questions → full README (up to KB_BUDGET_FULL_README chars)
      2. Fallback to raw-file RAG if README is absent or too thin (<200 chars
         of hand-written content outside the AUTO-INDEX block)

    Args:
        question:             The user's question
        folder_name:          KB folder to search (e.g. "ACE Docs")
        agent_name:           Display name for this agent
        system_prompt:        Domain-specific system prompt
        conversation_history: Optional prior messages for multi-turn context
        top_n:                Number of files to retrieve (fallback RAG only)
        max_chars:            Max chars to extract per file (fallback RAG only)
        _pre_ranked_results:  Optional pre-ordered search results supplied by a
                              domain-specific sub-agent.  When provided, skips
                              the vector search step and uses these results
                              directly.  This is the sub-agent hook: each
                              domain's agent_<safe>.py can do its own retrieval
                              (with domain-specific re-ranking, pinning, etc.)
                              and hand the sorted list to ask() for extraction
                              + LLM dispatch.

    Returns:
        { "agent", "answer", "sources", "found" }
    """
    # ── Strategy 1: README-first ──────────────────────────────────────────────
    readme_context, source_label = _get_readme_context(folder_name, question)

    if readme_context:
        is_complex = _is_complex_question(question)
        print(
            f"  [{agent_name}] README-first "
            f"({'full' if is_complex else 'index'} mode, "
            f"{len(readme_context):,} chars)"
        )

        # ── Passthrough: emit context for Bob's Claude to answer ──────────────
        if _PASSTHROUGH:
            print(f"  [{agent_name}] Passthrough mode — emitting context for Bob")
            return emit_passthrough(
                question     = question,
                context      = readme_context,
                system_prompt= system_prompt,
                agent_name   = agent_name,
                source_label = source_label,
            )

        messages = [{"role": "system", "content": system_prompt}]

        if conversation_history:
            messages.extend(conversation_history[-_cb.get("history"):])

        messages.append({
            "role": "user",
            "content": (
                f"Use the following knowledge base content to answer the question.\n\n"
                f"--- {source_label} ---\n{readme_context}\n---\n\n"
                f"Question: {question}"
            ),
        })

        answer = call_llm(messages)

        readme_sources = [{"name": source_label, "path": f"{folder_name}/README", "score": 1.0}]
        return {
            "agent":              agent_name,
            "answer":             answer,
            "sources":            readme_sources,
            "confidence_footer":  format_confidence_footer(readme_sources),
            "found":              True,
        }

    # ── Strategy 2: Raw-file RAG fallback ────────────────────────────────────
    print(f"  [{agent_name}] Falling back to raw-file RAG (no usable README)")

    # A sub-agent (agents/agent_<safe>.py) may supply pre-ranked results so
    # its domain-specific retrieval logic (pinning, re-ranking) runs before
    # this function is called.  Use them directly; skip the vector search.
    if _pre_ranked_results is not None:
        results = _pre_ranked_results
    else:
        sys.path.insert(0, str(pathlib.Path(__file__).parent))
        from embeddings import search
        results = search(question, folder_name, top_n=top_n)

    if not results:
        return {
            "agent":   agent_name,
            "answer":  f"I could not find any relevant documents in {folder_name} to answer this question.",
            "sources": [],
            "found":   False,
        }

    # ── Revenue-file pin + boost ──────────────────────────────────────────────
    # For data/revenue questions the authoritative source is always
    # "*Revenue*.xlsx" files — NOT CRM/Won-Deals files.  Two steps:
    #
    # 1. PIN: scan the folder index for Revenue report files and inject them
    #    into the result list even if vector search ranked them outside top_n.
    #    This handles multi-product queries where the combined embedding drifts
    #    away from individual file summaries.
    #
    # 2. SORT: put all Revenue files first so the LLM reads the authoritative
    #    Rev Act @ PC numbers before any CRM Won column values.
    if _is_data_question(question):
        # Build a set of paths already in results
        result_paths = {r["path"] for r in results}

        # Load the folder index to find Revenue files not in the current results
        safe_name  = folder_to_safe_name(folder_name)
        idx_path   = VECTOR_STORE / f"{safe_name}_index.json"
        if idx_path.exists():
            try:
                idx_data = json.loads(idx_path.read_text())
                for entry in idx_data.get("entries", []):
                    name = entry.get("name", "")
                    if ("revenue" in name.lower() and name.lower().endswith(".xlsx")
                            and entry.get("path") not in result_paths):
                        # Pin this Revenue file at score=1.0 (pinned, not from search)
                        results.append({
                            "path":  entry["path"],
                            "name":  name,
                            "score": 1.0,
                            "folder": folder_name,
                        })
                        result_paths.add(entry["path"])
                        print(
                            f"  [{agent_name}] 📌 Pinned revenue file: {name}",
                            flush=True,
                        )
            except Exception:
                pass  # index unreadable → fall through with original results

        # Sort: Revenue Report files first, everything else preserves its order
        def _revenue_priority(r: dict) -> int:
            name = r.get("name", "").lower()
            return 0 if ("revenue" in name and name.endswith(".xlsx")) else 1
        results = sorted(results, key=_revenue_priority)

    context_blocks = []
    sources        = []
    for r in results:
        file_path = KB_ROOT / r["path"]
        if not file_path.exists():
            print(
                f"  [{agent_name}] ⚠ Indexed file no longer on disk, skipping: "
                f"{r['path']}  (re-run generate.py to update the index)",
                flush=True,
            )
            continue
        text = extract_full_text(file_path, max_chars=max_chars)
        context_blocks.append(
            f"--- Source: {r['name']} (relevance: {r['score']:.2f}) ---\n{text}"
        )
        sources.append({"name": r["name"], "path": r["path"], "score": r["score"]})

    if not sources:
        return {
            "agent":   agent_name,
            "answer":  f"Found index entries for {folder_name} but source files are missing.",
            "sources": [],
            "found":   False,
        }

    context = "\n\n".join(context_blocks)

    # ── Passthrough: emit context for Bob's Claude to answer ──────────────────
    if _PASSTHROUGH:
        print(f"  [{agent_name}] Passthrough mode — emitting context for Bob")
        return emit_passthrough(
            question      = question,
            context       = context,
            system_prompt = system_prompt,
            agent_name    = agent_name,
            source_label  = ", ".join(s["name"] for s in sources),
        )

    messages = [{"role": "system", "content": system_prompt}]

    if conversation_history:
        messages.extend(conversation_history[-_cb.get("history"):])

    messages.append({
        "role": "user",
        "content": (
            f"Use the following documents to answer the question.\n\n"
            f"{context}\n\n---\nQuestion: {question}"
        ),
    })

    answer = call_llm(messages)

    return {
        "agent":             agent_name,
        "answer":            answer,
        "sources":           sources,
        "confidence_footer": format_confidence_footer(sources),
        "found":             True,
    }
